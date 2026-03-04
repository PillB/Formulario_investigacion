import pytest

from report import alerta_temprana
from report.alerta_temprana import build_alerta_temprana_ppt
from report_builder import CaseData


@pytest.mark.skipif(not alerta_temprana.PPTX_AVAILABLE, reason="python-pptx no disponible")
def test_report_generation_with_missing_data_uses_placeholders(tmp_path):
    data = CaseData.from_mapping(
        {
            "caso": {"id_caso": "2026-0001", "tipo_informe": "Alerta temprana"},
            "analisis": {},
            "productos": [],
            "clientes": [],
            "colaboradores": [],
            "reclamos": [],
            "riesgos": [],
            "operaciones": [],
            "encabezado": {},
            "responsables": [],
        }
    )
    output = tmp_path / "alerta_missing_data.pptx"

    path = build_alerta_temprana_ppt(data, output, llm_helper=None)

    assert path.exists()

    from pptx import Presentation

    deck = Presentation(path)
    all_text = " ".join(
        shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text")
    )
    assert "Recomendaciones" in all_text
    assert "N/A" in all_text
    assert len(deck.slides) == 2
