import os

import tkinter as tk
from tkinter import ttk

import pytest

import app as app_module
from report import alerta_temprana
from report.alerta_temprana import (_fit_text_to_box, _synthesize_section_text,
                                    build_alerta_temprana_ppt,
                                    SpanishSummaryHelper)
from report.alerta_temprana_content import (build_alerta_temprana_sections,
                                            build_executive_summary,
                                            _truncate, MAX_BULLETS,
                                            PLACEHOLDER)
from report_builder import CaseData


@pytest.mark.skipif(
    os.name != "nt" and not os.environ.get("DISPLAY"),
    reason="Tkinter no disponible en el entorno de pruebas",
)
def test_alerta_temprana_button_invokes_command(monkeypatch, messagebox_spy):
    try:
        root = tk.Tk()
        root.withdraw()
    except tk.TclError:
        pytest.skip("Tkinter no disponible en el entorno de pruebas")

    app = app_module.FraudCaseApp(root)
    invoked = []

    def fake_generate():
        invoked.append(True)

    monkeypatch.setattr(app, "generate_alerta_temprana_ppt", fake_generate)

    notebook = ttk.Notebook(root)
    actions_tab = ttk.Frame(notebook)
    app.build_actions_tab(actions_tab)

    app.btn_alerta_temprana.invoke()

    assert invoked
    root.destroy()


@pytest.mark.skipif(
    os.name != "nt" and not os.environ.get("DISPLAY"),
    reason="Tkinter no disponible en el entorno de pruebas",
)
def test_generate_alerta_temprana_ppt_uses_report_helper(monkeypatch, messagebox_spy):
    try:
        root = tk.Tk()
        root.withdraw()
    except tk.TclError:
        pytest.skip("Tkinter no disponible en el entorno de pruebas")

    app = app_module.FraudCaseApp(root)
    app.btn_alerta_temprana = ttk.Button(root)
    calls = {}

    def fake_generate_report_file(extension, builder, description, source_widget=None):
        calls["extension"] = extension
        calls["builder"] = builder
        calls["description"] = description
        calls["source_widget"] = source_widget

    monkeypatch.setattr(app, "_generate_report_file", fake_generate_report_file)

    app.generate_alerta_temprana_ppt()

    assert calls["extension"] == "pptx"
    assert calls["builder"] is build_alerta_temprana_ppt
    assert calls["source_widget"] is app.btn_alerta_temprana
    assert "Alerta temprana" in calls["description"]
    root.destroy()


@pytest.mark.skipif(not alerta_temprana.PPTX_AVAILABLE, reason="python-pptx no disponible")
def test_build_alerta_temprana_ppt_generates_presentation(tmp_path):
    class StubLLM:
        def __init__(self):
            self.prompts = []

        def summarize(self, section, prompt, *, max_new_tokens=None):
            self.prompts.append((section, prompt))
            return f"{section} sintetizado"

    data = CaseData.from_mapping(
        {
            "caso": {
                "id_caso": "2025-0001",
                "tipo_informe": "Fraude",
                "investigador_nombre": "Ana Investigadora",
                "categoria1": "Tarjeta",
                "modalidad": "Digital",
                "fecha_de_ocurrencia": "2025-01-01",
                "fecha_de_descubrimiento": "2025-01-02",
            },
            "clientes": [{"id_cliente": "CLI1"}],
            "colaboradores": [{"nombres": "Carlos", "flag": "involucrado", "area": "TI"}],
            "productos": [
                {
                    "monto_investigado": "100.00",
                    "monto_perdida_fraude": "10.00",
                    "monto_falla_procesos": "5.00",
                    "monto_contingencia": "3.00",
                    "monto_recuperado": "2.00",
                }
            ],
            "reclamos": [],
            "involucramientos": [],
            "riesgos": [{"id_riesgo": "R1"}],
            "normas": [],
            "analisis": {
                "antecedentes": {"text": "Se detectaron cargos inusuales en tarjetas digitales."},
                "recomendaciones": {"text": "Suspender temporalmente los accesos y revisar monitoreo."},
            },
            "encabezado": {},
            "operaciones": [],
            "anexos": [],
            "firmas": [],
            "recomendaciones_categorias": {},
        }
    )

    output = tmp_path / "alerta_temprana.pptx"
    stub_llm = StubLLM()
    path = build_alerta_temprana_ppt(data, output, llm_helper=stub_llm)
    assert path.exists()

    from pptx import Presentation

    deck = Presentation(path)
    all_text = " ".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert len(deck.slides) == 2
    assert "Alerta temprana" in all_text
    assert "Resumen ejecutivo" in all_text
    assert "2025-0001" in all_text
    assert "Cronología" in all_text
    assert "Riesgos identificados" in all_text
    assert "Recomendaciones" in all_text
    assert "sintetizado" in all_text
    assert any("Fraude" in prompt for _section, prompt in stub_llm.prompts)


def test_synthesize_section_text_uses_fallback_when_llm_missing():
    caso = {"id_caso": "123", "tipo_informe": "Fraude", "modalidad": "Digital"}
    analisis = {"comentario_breve": {"text": "Hubo transferencias sospechosas sin autorización."}}
    productos = [{"fecha_ocurrencia": "2025-01-01"}]
    riesgos = [{"id_riesgo": "R1", "descripcion": "Credenciales filtradas"}]
    operaciones = [{"accion": "Bloquear tarjeta", "estado": "Pendiente"}]
    colaboradores = [{"nombres": "Ana", "flag": "involucrado", "area": "Seguridad"}]

    class NullLLM(SpanishSummaryHelper):
        def summarize(self, section, prompt, *, max_new_tokens=None):
            return None

    sections = build_alerta_temprana_sections(
        {
            "caso": caso,
            "analisis": analisis,
            "productos": productos,
            "riesgos": riesgos,
            "operaciones": operaciones,
            "colaboradores": colaboradores,
            "encabezado": {},
            "clientes": [],
            "reclamos": [],
        }
    )
    resumen = _synthesize_section_text("Resumen", sections, caso, NullLLM())
    assert "transferencias sospechosas" in resumen


def test_sections_empty_return_placeholder():
    sections = build_alerta_temprana_sections(
        {
            "caso": {},
            "analisis": {},
            "productos": [],
            "riesgos": [],
            "operaciones": [],
            "colaboradores": [],
            "encabezado": {},
            "clientes": [],
            "reclamos": [],
        }
    )
    assert "Mensaje clave" in sections["resumen"]
    assert "Puntos de soporte" in sections["resumen"]
    assert "Evidencias" in sections["resumen"]
    assert "N/A" in sections["resumen"]
    assert sections["cronologia"] == PLACEHOLDER
    assert sections["analisis"] == PLACEHOLDER


def test_resumen_section_includes_blocks_and_references():
    sections = build_alerta_temprana_sections(
        {
            "caso": {
                "fecha_de_ocurrencia": "2025-01-01",
                "fecha_de_descubrimiento": "2025-01-02",
            },
            "analisis": {"comentario_breve": "Mensaje principal del caso."},
            "productos": [
                {
                    "monto_investigado": "100.00",
                    "monto_perdida_fraude": "10.00",
                    "monto_falla_procesos": "5.00",
                    "monto_contingencia": "3.00",
                    "monto_recuperado": "2.00",
                }
            ],
            "riesgos": [],
            "operaciones": [],
            "colaboradores": [],
            "encabezado": {"dirigido_a": "Comité", "area_reporte": "Riesgos"},
            "clientes": [{"id_cliente": "CLI-1"}],
            "reclamos": [],
        }
    )
    resumen = sections["resumen"]
    assert "Mensaje clave" in resumen
    assert "Puntos de soporte" in resumen
    assert "Evidencias" in resumen
    assert "Clientes vinculados: 1" in resumen
    assert "[Caso: fecha_de_ocurrencia]" in resumen


def test_recomendaciones_section_prefers_analisis_recomendaciones():
    sections = build_alerta_temprana_sections(
        {
            "caso": {"id_caso": "2025-0004"},
            "analisis": {
                "recomendaciones": "• Revisar límites de autorización.",
                "acciones": "• Esta acción no debería usarse.",
            },
            "productos": [],
            "riesgos": [],
            "operaciones": [],
            "colaboradores": [],
            "encabezado": {},
            "clientes": [],
            "reclamos": [],
        }
    )
    assert "Revisar límites de autorización" in sections["recomendaciones"]
    assert "no debería" not in sections["recomendaciones"]
    assert sections["acciones"] == sections["recomendaciones"]


def test_synthesize_section_text_skips_llm_when_sources_empty():
    class StubLLM:
        def __init__(self):
            self.prompts = []

        def summarize(self, section, prompt, *, max_new_tokens=None):
            self.prompts.append((section, prompt))
            return "contenido inventado"

    sections = build_alerta_temprana_sections(
        {
            "caso": {},
            "analisis": {},
            "productos": [],
            "riesgos": [],
            "operaciones": [],
            "colaboradores": [],
            "encabezado": {},
            "clientes": [],
            "reclamos": [],
        }
    )
    stub_llm = StubLLM()
    texto = _synthesize_section_text("Resumen", sections, {}, stub_llm)
    assert texto == PLACEHOLDER
    assert stub_llm.prompts == []


def test_cronologia_prefers_hallazgos_bullets():
    sections = build_alerta_temprana_sections(
        {
            "caso": {"id_caso": "2025-0002"},
            "analisis": {
                "hallazgos": (
                    "- Primer hallazgo relevante del caso.\n"
                    "- Segundo hallazgo con detalle adicional para cronología.\n"
                    "- Tercer hallazgo asociado a la investigación.\n"
                    "- Cuarto hallazgo operativo.\n"
                    "- Quinto hallazgo que no debe aparecer."
                )
            },
            "productos": [],
            "riesgos": [],
            "operaciones": [{"accion": "No debería usarse", "estado": "Pendiente"}],
            "colaboradores": [],
            "encabezado": {},
            "clientes": [],
            "reclamos": [],
        }
    )
    cronologia = sections["cronologia"]
    assert "Primer hallazgo relevante" in cronologia
    assert "Cuarto hallazgo operativo" in cronologia
    assert "Quinto hallazgo" not in cronologia


def test_cronologia_orders_operaciones_by_fecha_with_fallback():
    sections = build_alerta_temprana_sections(
        {
            "caso": {
                "id_caso": "2025-0012",
                "fecha_de_ocurrencia": "2025-01-05",
                "fecha_de_descubrimiento": "2025-01-10",
            },
            "analisis": {},
            "productos": [{"fecha_ocurrencia": "2025-01-07"}],
            "riesgos": [],
            "operaciones": [
                {"accion": "Cierre de investigación", "estado": "Completado", "fecha": "2025-03-05"},
                {"accion": "Revisión inicial", "estado": "Pendiente"},
                {"accion": "Entrevista", "estado": "En curso", "fecha": "2025-02-10"},
            ],
            "colaboradores": [],
            "encabezado": {},
            "clientes": [],
            "reclamos": [],
        }
    )
    cronologia_lines = [line.replace("•", "").strip() for line in sections["cronologia"].splitlines()]
    assert cronologia_lines[0].startswith("2025-01-05")
    assert cronologia_lines[1].startswith("2025-02-10")
    assert cronologia_lines[2].startswith("2025-03-05")
    assert "Revisión inicial" in cronologia_lines[0]
    assert "Entrevista" in cronologia_lines[1]
    assert "Cierre de investigación" in cronologia_lines[2]


def test_analisis_section_prioritizes_hallazgo_principal_and_control_failure():
    sections = build_alerta_temprana_sections(
        {
            "caso": {"id_caso": "2025-0010"},
            "analisis": {
                "hallazgos": "- Primer hallazgo crítico.\n- Segundo hallazgo adicional.",
                "conclusiones": "Falló el control de autenticación durante el proceso.",
                "antecedentes": "Se observaron alertas en el monitoreo.",
                "modus_operandi": "Uso de credenciales filtradas.",
            },
            "productos": [],
            "riesgos": [],
            "operaciones": [],
            "colaboradores": [],
            "encabezado": {},
            "clientes": [],
            "reclamos": [],
        }
    )
    analisis = sections["analisis"].splitlines()
    assert "Hallazgo principal" in analisis[0]
    assert "Fallo de control" in analisis[1]
    assert "Antecedentes" in analisis[2]
    assert "Modus operandi" in analisis[3]


def test_analisis_section_prioritizes_control_failure_without_hallazgos():
    sections = build_alerta_temprana_sections(
        {
            "caso": {"id_caso": "2025-0011"},
            "analisis": {
                "comentario_amplio": "Se verificó el expediente. Fallo de control en la validación.",
                "antecedentes": "Operación usual para el cliente.",
            },
            "productos": [],
            "riesgos": [],
            "operaciones": [],
            "colaboradores": [],
            "encabezado": {},
            "clientes": [],
            "reclamos": [],
        }
    )
    analisis = sections["analisis"].splitlines()
    assert "Fallo de control" in analisis[0]
    assert "Antecedentes" in analisis[1]


def test_sections_limit_and_truncate_bullets():
    long_detail = "Detalle " + ("x" * 220)
    data = {
        "caso": {"id_caso": "2025-0003", "investigador_nombre": "Equipo"},
        "analisis": {},
        "productos": [],
        "riesgos": [
            {"id_riesgo": f"R{i}", "descripcion": long_detail, "criticidad": "Alta"}
            for i in range(MAX_BULLETS + 2)
        ],
        "operaciones": [
            {"accion": long_detail, "cliente": "Cliente", "estado": "Pendiente"}
            for _ in range(MAX_BULLETS + 1)
        ],
        "colaboradores": [
            {"nombres": f"Persona {i}", "flag": "involucrado", "area": long_detail}
            for i in range(MAX_BULLETS + 1)
        ],
        "encabezado": {},
        "clientes": [],
        "reclamos": [],
    }
    sections = build_alerta_temprana_sections(data)
    riesgos_lines = sections["riesgos"].splitlines()
    recomendaciones_lines = sections["recomendaciones"].splitlines()
    responsables_lines = sections["responsables"].splitlines()

    assert len(riesgos_lines) == MAX_BULLETS
    assert len(recomendaciones_lines) == MAX_BULLETS
    assert len(responsables_lines) == MAX_BULLETS
    assert any(("…" in line) or ("[Riesgos: Registro]" in line) for line in riesgos_lines)
    assert any(("…" in line) or ("[Operaciones: Accion]" in line) for line in recomendaciones_lines)
    assert any("…" in line for line in responsables_lines)




def test_sections_include_traceability_references_in_key_outputs():
    sections = build_alerta_temprana_sections(
        {
            "caso": {
                "fecha_de_ocurrencia": "2025-01-01",
                "fecha_de_descubrimiento": "2025-01-02",
            },
            "analisis": {
                "comentario_breve": "Mensaje principal del caso.",
                "hallazgos": "- Hallazgo crítico.",
                "recomendaciones": "- Reforzar controles de autenticación.",
            },
            "productos": [{"monto_investigado": "100.00"}],
            "riesgos": [{"id_riesgo": "R-1", "descripcion": "Acceso indebido", "criticidad": "Alta"}],
            "operaciones": [],
            "colaboradores": [],
            "encabezado": {"dirigido_a": "Comité", "area_reporte": "Riesgos"},
            "clientes": [{"id_cliente": "CLI-1"}],
            "reclamos": [],
        }
    )

    assert "[Analisis: Comentario breve]" in sections["resumen"]
    assert "[Caso: fecha_de_ocurrencia]" in sections["resumen"]
    assert "[Analisis: Hallazgos]" in sections["analisis"]
    assert "[Riesgos: Registro]" in sections["riesgos"]
    assert "[Analisis: Recomendaciones]" in sections["recomendaciones"]


def test_truncation_preserves_source_reference_tokens():
    long_text = "Detalle " + ("x" * 260) + " [Analisis: Hallazgos]"
    result = _truncate(long_text, 120, label="prueba_referencia")
    assert "[Analisis: Hallazgos]" in result
    assert len(result) <= 120


def test_executive_summary_evidence_contains_source_references():
    summary = build_executive_summary(
        {
            "caso": {
                "id_caso": "2025-0450",
                "fecha_de_ocurrencia": "2025-03-01",
                "fecha_de_descubrimiento": "2025-03-03",
                "categoria1": "Tarjeta",
                "modalidad": "Digital",
            },
            "encabezado": {"dirigido_a": "Comité", "area_reporte": "Riesgos"},
            "productos": [{"id_producto": "P-1"}],
            "clientes": [{"id_cliente": "C-1"}],
            "colaboradores": [{"nombres": "Ana"}],
            "riesgos": [{"id_riesgo": "R-1"}],
            "reclamos": [{"id_reclamo": "CLM-1"}],
            "analisis": {},
            "operaciones": [],
            "responsables": [],
        }
    )

    assert any("[Caso: fecha_de_ocurrencia]" in line for line in summary.evidence)
    assert any("[Productos: Tabla]" in line for line in summary.evidence)
def test_fit_text_to_box_truncates_long_narrative(caplog):
    long_text = ("Narrativa extensa " * 400).strip()
    with caplog.at_level("WARNING"):
        result = _fit_text_to_box(
            long_text,
            width_in=3.5,
            height_in=1.2,
            section_title="Cronología",
        )
    assert result.truncated is True
    assert result.text.endswith("…")
    assert "Cronología" in caplog.text


@pytest.mark.skipif(not alerta_temprana.PPTX_AVAILABLE, reason="python-pptx no disponible")
def test_add_section_panel_supports_bullet_paragraphs():
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from pptx.util import Inches

    deck = Presentation()
    deck.slide_width = Inches(13.33)
    deck.slide_height = Inches(7.5)
    slide = deck.slides.add_slide(deck.slide_layouts[6])

    alerta_temprana._add_section_panel(
        slide,
        Inches(0.5),
        Inches(0.5),
        Inches(5.0),
        Inches(2.5),
        "Detalle",
        "• Primer punto\nSegundo párrafo\n- Tercer punto",
    )

    body_box = slide.shapes[-1]
    paragraphs = body_box.text_frame.paragraphs
    assert [para.text for para in paragraphs[:3]] == ["Primer punto", "Segundo párrafo", "Tercer punto"]

    first_ppr = paragraphs[0]._p.get_or_add_pPr()
    second_ppr = paragraphs[1]._p.get_or_add_pPr()
    third_ppr = paragraphs[2]._p.get_or_add_pPr()
    assert first_ppr.find(qn("a:buChar")) is not None
    assert second_ppr.find(qn("a:buChar")) is None
    assert third_ppr.find(qn("a:buChar")) is not None


def test_responsables_section_prioritizes_explicit_roles_over_colaboradores():
    sections = build_alerta_temprana_sections(
        {
            "caso": {"id_caso": "2025-0200", "investigador_nombre": "Líder"},
            "analisis": {},
            "productos": [{"id_producto": "P-001"}],
            "riesgos": [],
            "operaciones": [],
            "colaboradores": [
                {"nombres": "Persona Involucrada", "flag": "Involucrado", "area": "Backoffice"}
            ],
            "responsables": [
                {
                    "scope": "unidad",
                    "nombre": "Ana Unidad",
                    "puesto": "Jefa",
                    "division": "DCA",
                    "area": "Área comercial",
                    "servicio": "Atención",
                    "nombre_agencia": "Agencia Centro",
                },
                {
                    "scope": "producto",
                    "nombre": "Luis Producto",
                    "puesto": "Owner",
                    "id_producto": "P-001",
                },
            ],
            "encabezado": {},
            "clientes": [],
            "reclamos": [],
        }
    )

    responsables = sections["responsables"]
    assert "Responsable de unidad: Ana Unidad" in responsables
    assert "Responsable de producto: Luis Producto" in responsables
    assert "Persona Involucrada" not in responsables


def test_responsables_section_falls_back_to_colaboradores_when_roles_missing():
    sections = build_alerta_temprana_sections(
        {
            "caso": {"id_caso": "2025-0201"},
            "analisis": {},
            "productos": [],
            "riesgos": [],
            "operaciones": [],
            "colaboradores": [
                {"nombres": "Persona Involucrada", "flag": "Relacionado", "area": "Canales"}
            ],
            "responsables": [],
            "encabezado": {},
            "clientes": [],
            "reclamos": [],
        }
    )

    assert "Persona Involucrada (Relacionado - Canales)" in sections["responsables"]
