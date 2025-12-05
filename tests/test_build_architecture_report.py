from pathlib import Path

from PIL import Image
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

import pytest

import build_architecture_report as arch_report
from reportlab.platypus import Spacer


def _make_png(path: Path, size: tuple[int, int]) -> None:
    image = Image.new("RGB", size, color="white")
    image.save(path)


def test_render_mermaid_uses_high_resolution(monkeypatch, tmp_path):
    source = tmp_path / "diagram.mmd"
    target = tmp_path / "diagram.png"
    source.write_text("graph TD; A-->B;")

    monkeypatch.setattr(arch_report.shutil, "which", lambda name: "/usr/bin/mmdc" if name == "mmdc" else None)
    commands: list[list[str]] = []

    def fake_run(cmd: list[str], check: bool):  # noqa: D401
        commands.append(cmd)

    monkeypatch.setattr(arch_report.subprocess, "run", fake_run)

    arch_report.render_mermaid(source, target)

    assert commands, "render_mermaid should invoke mmdc"
    cmd = commands[0]
    assert "-s" in cmd
    assert str(arch_report.MERMAID_EXPORT_SCALE) in cmd
    assert "-w" in cmd
    assert str(arch_report.MERMAID_EXPORT_WIDTH_PX) in cmd


def test_build_stylesheet_reuses_existing_heading_styles():
    first = arch_report._build_stylesheet()
    assert first["Heading1"].fontSize == 18
    assert first["Heading2"].fontSize == 14

    second = arch_report._build_stylesheet()
    assert second["Heading1"].fontSize == 18
    assert second["Heading2"].fontSize == 14


def test_flowable_image_preserves_aspect_ratio(monkeypatch, tmp_path):
    image_path = tmp_path / "diagram.png"
    image_path.touch()

    def fake_reader(path: str):
        assert path == str(image_path)

        class Reader:
            def getSize(self):
                return (800, 400)

        return Reader()

    captured: dict[str, float] = {}

    def fake_image(path: str, width: float | None = None, height: float | None = None):
        captured["path"] = path
        captured["width"] = width
        captured["height"] = height
        return f"image:{path}"

    monkeypatch.setattr(arch_report, "ImageReader", fake_reader)
    monkeypatch.setattr(arch_report, "Image", fake_image)

    flowable = arch_report._flowable_image(image_path, target_width=6.5 * arch_report.inch)

    assert flowable == f"image:{image_path}"
    assert captured["path"] == str(image_path)
    assert captured["width"] == 6.5 * arch_report.inch
    assert captured["height"] == (6.5 * arch_report.inch) * 0.5


def test_build_report_avoids_duplicate_styles(monkeypatch, tmp_path):
    def fake_render(source: Path, target: Path) -> Path:
        target.parent.mkdir(parents=True, exist_ok=True)
        target.touch()
        return target

    monkeypatch.setattr(arch_report, "render_mermaid", fake_render)
    monkeypatch.setattr(arch_report, "_flowable_image", lambda *args, **kwargs: Spacer(1, 1))

    output = tmp_path / "architecture.pdf"
    generated = arch_report.build_report(output)

    assert generated == output
    assert output.exists()


def test_build_report_expands_diagram_pages(monkeypatch, tmp_path):
    def fake_render(source: Path, target: Path) -> Path:
        target.parent.mkdir(parents=True, exist_ok=True)
        target.touch()
        return target

    image_calls: list[tuple[Path, float]] = []

    def fake_flowable(image_path: Path, target_width: float):
        image_calls.append((image_path, target_width))
        return Spacer(1, 1)

    monkeypatch.setattr(arch_report, "render_mermaid", fake_render)
    monkeypatch.setattr(arch_report, "_flowable_image", fake_flowable)

    output = tmp_path / "architecture.pdf"
    arch_report.build_report(output)

    assert [call[0] for call in image_calls] == [arch_report.ARCH_PNG, arch_report.SEQ_PNG]

    arch_page_size = arch_report.landscape(arch_report.A3)
    expected_arch_width = arch_page_size[0] - (0.7 * arch_report.inch * 2)
    seq_page_size = arch_report.A3
    expected_seq_width = seq_page_size[0] - (0.7 * arch_report.inch * 2)

    assert image_calls[0][1] == pytest.approx(expected_arch_width)
    assert image_calls[1][1] == pytest.approx(expected_seq_width)


def test_build_editable_deck_uses_native_shapes(tmp_path):
    output = tmp_path / "diagrams.pptx"
    deck_path = arch_report.build_editable_deck(output)

    presentation = Presentation(deck_path)
    assert presentation.slide_width == arch_report.PPTX_SLIDE_WIDTH
    assert presentation.slide_height == arch_report.PPTX_SLIDE_HEIGHT

    assert len(presentation.slides) == 2
    arch_slide = presentation.slides[0]
    seq_slide = presentation.slides[1]

    assert all(
        shape.shape_type != MSO_SHAPE_TYPE.PICTURE for shape in arch_slide.shapes
    ), "Architecture slide must use native shapes"

    connector_names = {
        shape.name for shape in arch_slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.LINE
    }
    expected_import_links = {
        "importer_to_caseTab",
        "importer_to_clientTab",
        "importer_to_teamTab",
        "importer_to_productTab",
        "importer_to_riskTab",
        "importer_to_normTab",
        "importer_to_summary",
    }
    assert expected_import_links.issubset(connector_names)

    text_content = " ".join(
        shape.text for shape in arch_slide.shapes if hasattr(shape, "text")
    )
    assert "Servicio de Catálogos" in text_content
    assert "Pestaña Team Members" in text_content

    seq_text = " ".join(
        shape.text for shape in seq_slide.shapes if hasattr(shape, "text")
    )
    assert "Validación de integridad referencial" in seq_text
    assert "Autocompletado/Autopoblado" in seq_text


def test_architecture_layout_avoids_overlap_and_uses_black_font(tmp_path):
    deck_path = arch_report.build_editable_deck(tmp_path / "grid.pptx")
    presentation = Presentation(deck_path)
    arch_slide = presentation.slides[0]

    diagram_shapes = [
        shape
        for shape in arch_slide.shapes
        if shape.shape_type in {MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.TEXT_BOX}
    ]

    def overlaps(a: tuple[int, int, int, int], b: tuple[int, int, int, int]) -> bool:
        return not (a[2] <= b[0] or a[0] >= b[2] or a[3] <= b[1] or a[1] >= b[3])

    rectangles: list[tuple[int, int, int, int]] = []
    for shape in diagram_shapes:
        rect = (shape.left, shape.top, shape.left + shape.width, shape.top + shape.height)
        for other in rectangles:
            assert not overlaps(rect, other), f"{shape.name} overlaps another shape"
        rectangles.append(rect)

    appcore = next(shape for shape in diagram_shapes if shape.name == "container_appcore")
    assert appcore.text_frame.paragraphs[0].font.color.rgb == arch_report.PPTX_FONT_COLOR

    importer_link = next(
        shape
        for shape in arch_slide.shapes
        if getattr(shape, "name", "") == "importer_to_caseTab"
    )
    connector_format = getattr(importer_link, "connector_format", None)
    if connector_format:
        assert connector_format.type == arch_report.MSO_CONNECTOR.ELBOW
    else:
        assert importer_link.shape_type == MSO_SHAPE_TYPE.LINE


def test_main_generates_pptx_by_default(monkeypatch, tmp_path):
    pdf_out = tmp_path / "out.pdf"
    pptx_out = tmp_path / "out.pptx"

    calls: dict[str, Path] = {}

    def fake_build_report(path: Path):
        path.touch()
        calls["pdf"] = path
        return path

    def fake_build_pptx(path: Path):
        path.touch()
        calls["pptx"] = path
        return path

    monkeypatch.setattr(arch_report, "DEFAULT_OUTPUT", pdf_out)
    monkeypatch.setattr(arch_report, "DEFAULT_PPTX", pptx_out)
    monkeypatch.setattr(arch_report, "build_report", fake_build_report)
    monkeypatch.setattr(arch_report, "build_editable_deck", fake_build_pptx)

    exit_code = arch_report.main([])

    assert exit_code == 0
    assert calls["pdf"] == pdf_out
    assert calls["pptx"] == pptx_out


def test_main_can_skip_pptx(monkeypatch, tmp_path):
    pdf_out = tmp_path / "only.pdf"

    calls: dict[str, Path] = {}

    def fake_build_report(path: Path):
        path.touch()
        calls["pdf"] = path
        return path

    def fake_build_pptx(path: Path):  # pragma: no cover - should not be hit
        raise AssertionError("pptx should be skipped")

    monkeypatch.setattr(arch_report, "build_report", fake_build_report)
    monkeypatch.setattr(arch_report, "build_editable_deck", fake_build_pptx)

    exit_code = arch_report.main(["--output", str(pdf_out), "--no-pptx"])

    assert exit_code == 0
    assert calls["pdf"] == pdf_out
    assert "pptx" not in calls


def test_architecture_diagram_links_imports_and_autofill():
    content = Path(__file__).resolve().parent.parent.joinpath("docs", "architecture.mmd").read_text(
        encoding="utf-8"
    )

    assert "Aplicación Python Tkinter GUI" in content
    assert "FormularioInvestigacionesFraude" in content
    assert "Pestañas de Caso, Clientes, Team Members, Productos, Riesgos, Normas" in content
    assert "Tabs de Caso" not in content
    assert "FraudCaseApp" not in content

    importer_targets = [
        "caseTab",
        "clientTab",
        "teamTab",
        "productTab",
        "riskTab",
        "normTab",
        "summary",
    ]
    for target in importer_targets:
        assert f"importer --> {target}" in content

    autofill_targets = [
        "caseTab",
        "clientTab",
        "teamTab",
        "productTab",
        "riskTab",
        "normTab",
        "actions",
        "summary",
    ]
    for target in autofill_targets:
        assert f"{target} --> autofill" in content

    assert "\\n" not in content
    assert "Autofill" not in content
    assert "CatalogService" not in content


def test_sequence_diagram_is_spanish_and_covers_data_sources():
    content = Path(__file__).resolve().parent.parent.joinpath("docs", "sequence_diagram.mmd").read_text(
        encoding="utf-8"
    )

    assert "UI Python Tkinter GUI (FormularioInvestigacionesFraude)" in content
    assert "Servicio de Catálogos/Autopoblado" in content
    assert "Importador CSV" in content
    assert "Note over UI,Import: Carga masiva" in content
    assert "UI->>UI: _apply_*_import_payload" in content
    assert "Validación por fila y duplicados (cliente/colaborador/riesgo/producto/norma)" in content
    assert "Reporte: build_report + build_docx + build_editable_deck" in content
    assert "Note over UI,Catalogos: Autocompletado/Autopoblado" in content
    assert "Validación de integridad referencial (IDs existen en catálogos/snapshots)" in content
    assert "Autofill" not in content
    assert "FraudCaseApp" not in content
    assert "CatalogService" not in content
    assert "\\n" not in content
