"""Generate a pitch-style PowerPoint deck for the Formulario de Investigación app."""
from __future__ import annotations

import argparse
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from openpyxl import load_workbook
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

PROJECT_ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = PROJECT_ROOT / "docs"
WIREFRAMES_DIR = PROJECT_ROOT / "wireframes"
WIRE_FRAMEBOOK = DOCS_DIR / "formulario_investigaciones_wireframe.xlsx"

DEFAULT_OUTPUT = PROJECT_ROOT / "Formulario_Investigacion_Pitch_Deck.pptx"
ASSETS_DIR = DOCS_DIR / "pitch_deck_assets"

SLIDE_WIDTH = Cm(33.867)  # 16:9 widescreen
SLIDE_HEIGHT = Cm(19.05)
MARGIN_X = Cm(1.2)
MARGIN_Y = Cm(0.8)
TITLE_HEIGHT = Cm(1.6)
SUBTITLE_HEIGHT = Cm(0.9)
CONTENT_GAP = Cm(0.4)

TITLE_SIZE = Pt(28)
SUBTITLE_SIZE = Pt(16)
BODY_SIZE = Pt(14)
BODY_SMALL_SIZE = Pt(12)

COLOR_DARK = RGBColor(20, 28, 38)
COLOR_ACCENT = RGBColor(32, 90, 166)
COLOR_SOFT = RGBColor(236, 243, 252)
COLOR_MUTED = RGBColor(99, 115, 129)
COLOR_BORDER = RGBColor(204, 214, 221)


@dataclass(frozen=True)
class DiagramSpec:
    source: Path
    output: Path
    title: str


def render_mermaid(source: Path, target: Path) -> Path:
    """Render Mermaid to PNG using CLI or generate placeholder images."""
    import shutil
    import subprocess

    if not source.exists():
        raise FileNotFoundError(source)

    target.parent.mkdir(parents=True, exist_ok=True)
    cmd: list[str] | None
    if shutil.which("mmdc"):
        cmd = ["mmdc", "-i", str(source), "-o", str(target)]
    elif shutil.which("npx"):
        cmd = [
            "npx",
            "-y",
            "@mermaid-js/mermaid-cli",
            "-i",
            str(source),
            "-o",
            str(target),
        ]
    else:
        return render_placeholder(source, target, "Mermaid CLI no disponible")

    cmd.extend(["-s", "2.3", "-w", "2200"])
    if (PROJECT_ROOT / "puppeteer-config.json").exists():
        cmd.extend(["-p", str(PROJECT_ROOT / "puppeteer-config.json")])
    if (PROJECT_ROOT / "mermaid-config.json").exists():
        cmd.extend(["-c", str(PROJECT_ROOT / "mermaid-config.json")])

    try:
        subprocess.run(cmd, check=True)
    except subprocess.CalledProcessError:
        return render_placeholder(source, target, "Fallo al ejecutar Mermaid CLI")
    return target


def render_placeholder(source: Path, target: Path, reason: str) -> Path:
    """Build a simple placeholder image when diagrams cannot be rendered."""
    target.parent.mkdir(parents=True, exist_ok=True)
    img = Image.new("RGB", (2200, 1400), color="white")
    draw = ImageDraw.Draw(img)
    draw.rectangle((30, 30, 2170, 1370), outline="black", width=6)
    text = f"{source.name}\n{reason}\nInstala mermaid-cli para renderizar el diagrama"
    draw.multiline_text((80, 80), text, fill="black", font=ImageFont.load_default(), spacing=10)
    img.save(target)
    return target


def add_title_block(slide, title: str, subtitle: str | None = None) -> float:
    left = MARGIN_X
    top = MARGIN_Y
    width = SLIDE_WIDTH - 2 * MARGIN_X
    title_box = slide.shapes.add_textbox(left, top, width, TITLE_HEIGHT)
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = TITLE_SIZE
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = COLOR_DARK

    content_top = top + TITLE_HEIGHT
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(left, content_top, width, SUBTITLE_HEIGHT)
        subtitle_tf = subtitle_box.text_frame
        subtitle_tf.text = subtitle
        subtitle_tf.paragraphs[0].font.size = SUBTITLE_SIZE
        subtitle_tf.paragraphs[0].font.color.rgb = COLOR_MUTED
        content_top += SUBTITLE_HEIGHT
    return content_top + CONTENT_GAP


def add_bullets(slide, items: Iterable[str], left: float, top: float, width: float, height: float) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.text = item
        paragraph.font.size = BODY_SIZE
        paragraph.font.color.rgb = COLOR_DARK
        paragraph.level = 0


def add_caption(slide, text: str, left: float, top: float, width: float, height: float) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.text = text
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.font.size = BODY_SMALL_SIZE
    para.font.color.rgb = COLOR_MUTED


def add_image(slide, image_path: Path, left: float, top: float, width: float, height: float) -> None:
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
        return
    placeholder = render_placeholder(image_path, ASSETS_DIR / f"missing_{image_path.stem}.png", "Imagen no encontrada")
    slide.shapes.add_picture(str(placeholder), left, top, width=width, height=height)


def load_wireframe_bullets(sheet_name: str, max_items: int = 7) -> list[str]:
    workbook = load_workbook(WIRE_FRAMEBOOK, read_only=True, data_only=True)
    ws = workbook[sheet_name]
    bullets: list[str] = []
    for row in ws.iter_rows(values_only=True):
        field, field_type, description = (row + (None, None, None))[:3]
        if not field:
            continue
        if field_type is None and description is None:
            continue
        summary = f"{field} ({field_type})" if field_type else str(field)
        if description:
            summary = f"{summary}: {description}"
        bullets.append(summary)
        if len(bullets) >= max_items:
            break
    return bullets


def build_slide_definition(
    title: str,
    subtitle: str,
    diagram: Path | None,
    bullets: list[str],
) -> tuple[str, str, Path | None, list[str]]:
    return title, subtitle, diagram, bullets


def build_diagram_assets() -> dict[str, Path]:
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    diagrams = {
        "sequence": DiagramSpec(DOCS_DIR / "sequence_diagram.mmd", ASSETS_DIR / "sequence.png", "Secuencia"),
        "architecture": DiagramSpec(DOCS_DIR / "architecture.mmd", ASSETS_DIR / "architecture.png", "Arquitectura"),
        "db": DiagramSpec(DOCS_DIR / "db_architecture.mmd", ASSETS_DIR / "db_architecture.png", "BD"),
        "circular": DiagramSpec(
            DOCS_DIR / "circular_eventos_flow.mmd", ASSETS_DIR / "circular_eventos_flow.png", "Eventos"
        ),
        "tabs_flow": DiagramSpec(
            WIREFRAMES_DIR / "layout_hierarchy.mmd", ASSETS_DIR / "layout_hierarchy.png", "Flujo UI"
        ),
        "tab_case": DiagramSpec(
            WIREFRAMES_DIR / "tab01_caso_participantes.mmd", ASSETS_DIR / "tab01_caso.png", "Caso"
        ),
        "tab_riesgos": DiagramSpec(
            WIREFRAMES_DIR / "tab02_riesgos.mmd", ASSETS_DIR / "tab02_riesgos.png", "Riesgos"
        ),
        "tab_normas": DiagramSpec(
            WIREFRAMES_DIR / "tab03_normas.mmd", ASSETS_DIR / "tab03_normas.png", "Normas"
        ),
        "tab_analisis": DiagramSpec(
            WIREFRAMES_DIR / "tab04_analisis.mmd", ASSETS_DIR / "tab04_analisis.png", "Análisis"
        ),
        "tab_acciones": DiagramSpec(
            WIREFRAMES_DIR / "tab05_acciones.mmd", ASSETS_DIR / "tab05_acciones.png", "Acciones"
        ),
        "tab_resumen": DiagramSpec(
            WIREFRAMES_DIR / "tab06_resumen.mmd", ASSETS_DIR / "tab06_resumen.png", "Resumen"
        ),
        "tab_validacion": DiagramSpec(
            WIREFRAMES_DIR / "tab07_validacion.mmd", ASSETS_DIR / "tab07_validacion.png", "Validación"
        ),
    }
    results: dict[str, Path] = {}
    for key, spec in diagrams.items():
        results[key] = render_mermaid(spec.source, spec.output)
    return results


def add_splash_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_shape(
        1,
        0,
        0,
        SLIDE_WIDTH,
        SLIDE_HEIGHT,
    ).fill.solid()
    background = slide.shapes[0]
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_SOFT
    background.line.color.rgb = COLOR_SOFT

    title_box = slide.shapes.add_textbox(MARGIN_X, Cm(4.3), SLIDE_WIDTH - 2 * MARGIN_X, Cm(4))
    tf = title_box.text_frame
    tf.text = "Formulario de Investigación de Fraude"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_DARK
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(MARGIN_X, Cm(8.2), SLIDE_WIDTH - 2 * MARGIN_X, Cm(2))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.text = "Pitch visual: visión, flujo y arquitectura de la aplicación"
    subtitle_tf.paragraphs[0].font.size = Pt(18)
    subtitle_tf.paragraphs[0].font.color.rgb = COLOR_MUTED
    subtitle_tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_objective_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_top = add_title_block(
        slide,
        "Objetivo y valor generado",
        "La app reemplaza hojas sueltas con un flujo gobernado y trazable.",
    )
    bullets = [
        "Centralizar el registro de casos y participantes en una sola interfaz guiada por pestañas.",
        "Asegurar calidad de datos con validaciones cruzadas (fechas, montos, IDs, duplicados) alineadas al Design document CM.",
        "Eliminar inconsistencias de consolidar Excel dispersos mediante catálogos oficiales y autopoblado.",
        "Generar reportes ejecutivos y exportes normalizados (eventos.csv) listos para Bakehouse/PDA.",
        "Mantener trazabilidad con autosave, logs y versiones históricas por caso.",
    ]
    add_bullets(slide, bullets, MARGIN_X, content_top, SLIDE_WIDTH - 2 * MARGIN_X, Cm(10))


def add_key_functionalities_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_top = add_title_block(
        slide,
        "Funcionalidades clave y controles de consistencia",
        "Vista rápida de lo que el arquitecto debe validar.",
    )

    left_col_width = (SLIDE_WIDTH - 2 * MARGIN_X) * 0.55
    right_col_width = (SLIDE_WIDTH - 2 * MARGIN_X) * 0.42
    left = MARGIN_X
    right = MARGIN_X + left_col_width + Cm(0.6)

    left_items = [
        "Pestañas alineadas al ciclo del caso: Caso, Clientes, Colaboradores, Productos, Riesgos, Normas, Análisis, Acciones y Resumen.",
        "Validaciones en línea: formato AAAA-NNNN, IDs, fechas coherentes y montos con sumas exactas.",
        "Regla de llave técnica para evitar duplicados (caso + producto + cliente + colaborador + fecha + reclamo).",
        "Panel de validación con mensajes accionables y navegación directa al campo con error.",
    ]
    add_bullets(slide, left_items, left, content_top, left_col_width, Cm(10))

    right_items = [
        "Exportes automáticos: eventos.csv, analisis.csv, clientes.csv y datasets históricos h_*.csv.",
        "Reportes clave: Carta de inmediatez, Informe de Gerencia, Alerta temprana (PPT) y resumen ejecutivo.",
        "Consistencia frente a Excel: catálogos maestros, autopoblado y auditoría de logs.",
        "Integración con Bakehouse/PDA: eventos.csv mantiene estructura canónica y legacy.",
    ]
    add_bullets(slide, right_items, right, content_top, right_col_width, Cm(10))


def add_sequence_slide(prs: Presentation, diagram_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_top = add_title_block(
        slide,
        "Secuencia operativa: usuario → validaciones → reportes",
        "Basado en docs/sequence_diagram.mmd.",
    )
    image_height = SLIDE_HEIGHT - content_top - MARGIN_Y - Cm(1.3)
    add_image(slide, diagram_path, MARGIN_X, content_top, SLIDE_WIDTH - 2 * MARGIN_X, image_height)
    add_caption(
        slide,
        "Resumen: la UI valida, autoguarda, registra logs y desencadena exportes/alertas de forma orquestada.",
        MARGIN_X,
        SLIDE_HEIGHT - MARGIN_Y - Cm(1.0),
        SLIDE_WIDTH - 2 * MARGIN_X,
        Cm(0.8),
    )


def add_flow_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_top = add_title_block(
        slide,
        "Flujo de datos de punta a punta",
        "Entrada guiada → validación → exportes listos para Bakehouse.",
    )
    box_width = Cm(5.1)
    box_height = Cm(2.0)
    gap = Cm(0.8)
    start_left = MARGIN_X
    top = content_top + Cm(1)
    labels = [
        "Ingreso de caso\npor pestañas",
        "Validaciones\n+ autosave",
        "Autopoblado\ncatálogos",
        "Reportes &\nexports",
        "eventos.csv\nBakehouse",
    ]

    shapes = []
    for idx, label in enumerate(labels):
        left = start_left + idx * (box_width + gap)
        shape = slide.shapes.add_shape(1, left, top, box_width, box_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_SOFT
        shape.line.color.rgb = COLOR_BORDER
        tf = shape.text_frame
        tf.text = label
        tf.paragraphs[0].font.size = BODY_SMALL_SIZE
        tf.paragraphs[0].font.color.rgb = COLOR_DARK
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        shapes.append(shape)

    for idx in range(len(shapes) - 1):
        left = shapes[idx].left + box_width
        top_line = shapes[idx].top + box_height / 2
        line = slide.shapes.add_shape(20, left, top_line - Cm(0.1), gap, Cm(0.2))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_ACCENT
        line.line.color.rgb = COLOR_ACCENT

    add_caption(
        slide,
        "El flujo reduce fricción frente a Excel: datos consistentes, trazables y listos para ingestión.",
        MARGIN_X,
        top + box_height + Cm(0.6),
        SLIDE_WIDTH - 2 * MARGIN_X,
        Cm(1),
    )


def add_report_exports_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_top = add_title_block(
        slide,
        "Reportes y exportes críticos",
        "Salidas estandarizadas para gerencia, cumplimiento y Bakehouse/PDA.",
    )
    left_col_width = (SLIDE_WIDTH - 2 * MARGIN_X) * 0.48
    right_col_width = (SLIDE_WIDTH - 2 * MARGIN_X) * 0.48
    left = MARGIN_X
    right = MARGIN_X + left_col_width + Cm(0.6)
    left_items = [
        "eventos.csv: export canónico/legacy con combinación Producto × Reclamo × Involucramiento.",
        "analisis.csv: narrativa y hallazgos detallados.",
        "clientes.csv / colaboradores.csv / productos.csv: entidades normalizadas para consolidación.",
        "logs.csv: auditoría de validaciones, navegación e importaciones.",
    ]
    right_items = [
        "Carta de inmediatez (DOCX + CSV): notificación formal al colaborador.",
        "Informe de Gerencia: resumen ejecutivo con tablas y montos clave.",
        "Alerta temprana (PPTX): síntesis visual con riesgos y cronología.",
        "Históricos h_*.csv: consolidación con respaldo en unidad externa.",
    ]
    add_bullets(slide, left_items, left, content_top, left_col_width, Cm(9))
    add_bullets(slide, right_items, right, content_top, right_col_width, Cm(9))


def add_architecture_slide(prs: Presentation, diagram_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_top = add_title_block(
        slide,
        "Arquitectura de la app (contenedores y servicios)",
        "Estructura real del repositorio y dependencias de exportación.",
    )
    image_height = SLIDE_HEIGHT - content_top - MARGIN_Y - Cm(1.0)
    add_image(slide, diagram_path, MARGIN_X, content_top, SLIDE_WIDTH - 2 * MARGIN_X, image_height)
    add_caption(
        slide,
        "Ubica UI, validadores, reportes y archivos externos dentro del ecosistema de investigación.",
        MARGIN_X,
        SLIDE_HEIGHT - MARGIN_Y - Cm(0.8),
        SLIDE_WIDTH - 2 * MARGIN_X,
        Cm(0.7),
    )


def add_tabs_architecture_slide(prs: Presentation, diagram_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_top = add_title_block(
        slide,
        "Arquitectura de pestañas y secciones",
        "Wireframe estructural de la navegación principal.",
    )
    image_height = SLIDE_HEIGHT - content_top - MARGIN_Y - Cm(1.2)
    add_image(slide, diagram_path, MARGIN_X, content_top, SLIDE_WIDTH - 2 * MARGIN_X, image_height)
    add_caption(
        slide,
        "Define la jerarquía visual y el orden de captura para garantizar consistencia.",
        MARGIN_X,
        SLIDE_HEIGHT - MARGIN_Y - Cm(0.8),
        SLIDE_WIDTH - 2 * MARGIN_X,
        Cm(0.7),
    )


def add_db_architecture_slide(prs: Presentation, diagram_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_top = add_title_block(
        slide,
        "Arquitectura de datos (modelo lógico)",
        "Referenciado a db_architecture.mmd.",
    )
    image_height = SLIDE_HEIGHT - content_top - MARGIN_Y - Cm(1.2)
    add_image(slide, diagram_path, MARGIN_X, content_top, SLIDE_WIDTH - 2 * MARGIN_X, image_height)
    add_caption(
        slide,
        "Entidades clave: casos, clientes, colaboradores, productos, riesgos, normas y llave técnica.",
        MARGIN_X,
        SLIDE_HEIGHT - MARGIN_Y - Cm(0.8),
        SLIDE_WIDTH - 2 * MARGIN_X,
        Cm(0.7),
    )


def add_circular_flow_slide(prs: Presentation, diagram_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_top = add_title_block(
        slide,
        "Flujo circular: formulario ↔ eventos.csv",
        "Garantiza la consistencia del export crítico para Bakehouse/PDA.",
    )
    image_height = SLIDE_HEIGHT - content_top - MARGIN_Y - Cm(1.2)
    add_image(slide, diagram_path, MARGIN_X, content_top, SLIDE_WIDTH - 2 * MARGIN_X, image_height)
    add_caption(
        slide,
        "Cada fila de eventos.csv combina producto × reclamo × involucramiento y conserva trazabilidad.",
        MARGIN_X,
        SLIDE_HEIGHT - MARGIN_Y - Cm(0.8),
        SLIDE_WIDTH - 2 * MARGIN_X,
        Cm(0.7),
    )


def add_wireframe_slide(prs: Presentation, title: str, subtitle: str, diagram_path: Path, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_top = add_title_block(slide, title, subtitle)
    left_col_width = (SLIDE_WIDTH - 2 * MARGIN_X) * 0.6
    right_col_width = (SLIDE_WIDTH - 2 * MARGIN_X) * 0.38

    add_image(
        slide,
        diagram_path,
        MARGIN_X,
        content_top,
        left_col_width,
        SLIDE_HEIGHT - content_top - MARGIN_Y,
    )
    add_bullets(
        slide,
        bullets,
        MARGIN_X + left_col_width + Cm(0.4),
        content_top,
        right_col_width,
        SLIDE_HEIGHT - content_top - MARGIN_Y,
    )


def add_reference_slide(prs: Presentation, title: str, subtitle: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_top = add_title_block(slide, title, subtitle)
    add_bullets(slide, bullets, MARGIN_X, content_top, SLIDE_WIDTH - 2 * MARGIN_X, Cm(10))


def build_deck(output_path: Path = DEFAULT_OUTPUT) -> Path:
    diagram_paths = build_diagram_assets()
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    add_splash_slide(prs)
    add_objective_slide(prs)
    add_key_functionalities_slide(prs)
    add_sequence_slide(prs, diagram_paths["sequence"])
    add_flow_slide(prs)
    add_report_exports_slide(prs)
    add_tabs_architecture_slide(prs, diagram_paths["tabs_flow"])
    add_architecture_slide(prs, diagram_paths["architecture"])
    add_db_architecture_slide(prs, diagram_paths["db"])
    add_circular_flow_slide(prs, diagram_paths["circular"])

    wireframe_config = [
        build_slide_definition(
            "Wireframe: Caso y participantes",
            "Sección central con datos del caso, clientes, productos y colaboradores.",
            diagram_paths["tab_case"],
            load_wireframe_bullets("Caso y Participantes"),
        ),
        build_slide_definition(
            "Wireframe: Riesgos",
            "Registro de riesgos y planes de acción vinculados.",
            diagram_paths["tab_riesgos"],
            load_wireframe_bullets("Riesgos"),
        ),
        build_slide_definition(
            "Wireframe: Normas",
            "Normativa transgredida y fechas de vigencia.",
            diagram_paths["tab_normas"],
            load_wireframe_bullets("Normas"),
        ),
        build_slide_definition(
            "Wireframe: Análisis y narrativas",
            "Narrativa del caso y secciones extendidas del informe.",
            diagram_paths["tab_analisis"],
            load_wireframe_bullets("Análisis y Narrativas"),
        ),
        build_slide_definition(
            "Wireframe: Acciones",
            "Operaciones de carga masiva, catálogos y exportes.",
            diagram_paths["tab_acciones"],
            load_wireframe_bullets("Acciones"),
        ),
        build_slide_definition(
            "Wireframe: Resumen",
            "Vista consolidada para control previo a exportar.",
            diagram_paths["tab_resumen"],
            load_wireframe_bullets("Resumen"),
        ),
    ]

    for title, subtitle, diagram, bullets in wireframe_config:
        if diagram is None:
            add_reference_slide(prs, title, subtitle, bullets)
        else:
            add_wireframe_slide(prs, title, subtitle, diagram, bullets)

    reference_slides = [
        build_slide_definition(
            "Export: eventos.csv (estructura crítica)",
            "Detalle de columnas y origen por pestaña.",
            None,
            load_wireframe_bullets("Eventos_CSV", max_items=9),
        ),
        build_slide_definition(
            "Registro de Logs",
            "Auditoría de interacciones, validaciones e importaciones.",
            None,
            load_wireframe_bullets("Logs", max_items=8),
        ),
        build_slide_definition(
            "Reporte: Carta de inmediatez",
            "Campos clave para la notificación formal.",
            None,
            load_wireframe_bullets("Carta_inmediatez", max_items=8),
        ),
        build_slide_definition(
            "Reporte: Informe de Gerencia",
            "Secciones y contenido ejecutivo.",
            None,
            load_wireframe_bullets("Informe_Gerencia", max_items=7),
        ),
        build_slide_definition(
            "Reporte: Alerta temprana",
            "Resumen visual para stakeholders.",
            None,
            load_wireframe_bullets("Alerta_Temprana", max_items=7),
        ),
        build_slide_definition(
            "Panel de validación",
            "Reglas críticas de formato y consistencia.",
            None,
            load_wireframe_bullets("Panel_Validacion", max_items=7),
        ),
        build_slide_definition(
            "Mapeo de catálogos",
            "Fuentes de datos maestros y autopoblado.",
            None,
            load_wireframe_bullets("Mapeo Catalogos", max_items=7),
        ),
    ]

    for title, subtitle, _diagram, bullets in reference_slides:
        add_reference_slide(prs, title, subtitle, bullets)

    prs.save(output_path)
    return output_path


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--output",
        type=Path,
        default=DEFAULT_OUTPUT,
        help="Ruta del archivo PPTX a generar.",
    )
    args = parser.parse_args(argv)
    output_path = build_deck(args.output)
    print(f"Presentación generada en {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
