from __future__ import annotations

from dataclasses import dataclass, field
import logging
import math
import importlib.util
from pathlib import Path
import re
from typing import Mapping

from report.alerta_temprana_content import (
    ExecutiveSummary,
    build_alerta_temprana_sections,
    build_executive_summary,
)
from validators import sanitize_rich_text
from report_builder import CaseData

logger = logging.getLogger(__name__)

pptx_spec = importlib.util.find_spec("pptx")
if pptx_spec:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Inches, Pt
else:  # pragma: no cover - dependencia opcional
    Presentation = None
    RGBColor = None
    MSO_SHAPE = None
    PP_ALIGN = None
    qn = None
    OxmlElement = None
    Inches = None
    Pt = None

PPTX_AVAILABLE = Presentation is not None
PPTX_MISSING_MESSAGE = (
    "La dependencia opcional 'python-pptx' no está instalada. "
    "Ejecuta 'pip install python-pptx' para habilitar la alerta temprana en PPT."
)
PLACEHOLDER = "N/A"
DEFAULT_MODEL = "PlanTL-GOB-ES/flan-t5-base-spanish"
MAX_SECTION_CHARS = 600

TRANSFORMERS_AVAILABLE = importlib.util.find_spec("transformers") is not None

SLIDE_WIDTH_16_9 = Inches(13.33) if Inches else 0
SLIDE_HEIGHT_16_9 = Inches(7.5) if Inches else 0
MARGIN = Inches(0.4) if Inches else 0
COLUMN_GAP = Inches(0.3) if Inches else 0
MASTHEAD_HEIGHT = Inches(1.1) if Inches else 0
SECTION_HEADER_HEIGHT = Inches(0.28) if Inches else 0
PANEL_PADDING = Inches(0.18) if Inches else 0
BULLET_PREFIX = re.compile(r"^\s*(?:[-*•]|\d+[.)])\s+")
DEFAULT_BODY_FONT_PT = 11
MIN_BODY_FONT_PT = 9
AVG_CHAR_WIDTH_FACTOR = 0.52
DEFAULT_LINE_SPACING = 1.15


@dataclass(frozen=True)
class FitTextResult:
    text: str
    font_pt: int
    truncated: bool


def _emu_to_inches(value: float | int) -> float:
    if not Inches:
        return float(value)
    return float(value) / float(Inches(1))


def _estimate_chars_per_line(width_in: float, font_pt: int) -> int:
    if width_in <= 0 or font_pt <= 0:
        return 1
    width_pt = width_in * 72
    char_width = font_pt * AVG_CHAR_WIDTH_FACTOR
    estimated = int(width_pt / max(char_width, 1))
    return max(estimated, 15)


def _estimate_line_capacity(height_in: float, font_pt: int, line_spacing: float) -> int:
    if height_in <= 0 or font_pt <= 0:
        return 1
    line_height_in = (font_pt / 72) * line_spacing
    return max(int(height_in / max(line_height_in, 0.01)), 1)


def _estimate_text_lines(text: str, chars_per_line: int) -> int:
    if not text:
        return 1
    normalized = text.splitlines() or [text]
    total = 0
    for line in normalized:
        length = max(len(line), 1)
        total += max(1, math.ceil(length / max(chars_per_line, 1)))
    return total


def _truncate_text_to_fit(text: str, max_chars: int) -> str:
    cleaned = sanitize_rich_text(text, max_chars=None).strip()
    if max_chars <= 0 or len(cleaned) <= max_chars:
        return cleaned
    truncated = cleaned[: max_chars - 1].rstrip()
    last_space = truncated.rfind(" ")
    min_space_index = int(max_chars * 0.6)
    if last_space >= min_space_index:
        truncated = truncated[:last_space].rstrip()
    return truncated + "…"


def _fit_text_to_box(
    text: str,
    width_in: float,
    height_in: float,
    *,
    section_title: str,
    base_font_pt: int = DEFAULT_BODY_FONT_PT,
    min_font_pt: int = MIN_BODY_FONT_PT,
    line_spacing: float = DEFAULT_LINE_SPACING,
) -> FitTextResult:
    cleaned = sanitize_rich_text(text, max_chars=None).strip()
    if not cleaned:
        return FitTextResult(text=PLACEHOLDER, font_pt=base_font_pt, truncated=False)
    for font_pt in range(base_font_pt, min_font_pt - 1, -1):
        max_lines = _estimate_line_capacity(height_in, font_pt, line_spacing)
        chars_per_line = _estimate_chars_per_line(width_in, font_pt)
        total_lines = _estimate_text_lines(cleaned, chars_per_line)
        if total_lines <= max_lines:
            return FitTextResult(text=cleaned, font_pt=font_pt, truncated=False)
    max_lines = _estimate_line_capacity(height_in, min_font_pt, line_spacing)
    chars_per_line = _estimate_chars_per_line(width_in, min_font_pt)
    max_chars = max_lines * chars_per_line
    truncated = _truncate_text_to_fit(cleaned, max_chars)
    if truncated != cleaned:
        logger.warning(
            "Se truncó la sección '%s' para ajustarse al panel (%s → %s caracteres).",
            section_title,
            len(cleaned),
            len(truncated),
        )
    return FitTextResult(text=truncated, font_pt=min_font_pt, truncated=truncated != cleaned)


def _set_paragraph_bullet(paragraph, enabled: bool) -> None:
    if not enabled or not PPTX_AVAILABLE:
        return
    ppr = paragraph._p.get_or_add_pPr()
    existing = ppr.find(qn("a:buChar"))
    if existing is not None:
        ppr.remove(existing)
    bullet = OxmlElement("a:buChar")
    bullet.set("char", "•")
    ppr.insert(0, bullet)


def _split_body_paragraphs(body: str) -> list[tuple[str, bool]]:
    lines = [line.strip() for line in str(body or "").splitlines() if line.strip()]
    if not lines:
        return [(PLACEHOLDER, False)]
    parsed = []
    for line in lines:
        is_bullet = bool(BULLET_PREFIX.match(line))
        text = BULLET_PREFIX.sub("", line).strip() if is_bullet else line
        if text:
            parsed.append((text, is_bullet))
    return parsed or [(PLACEHOLDER, False)]


def _is_placeholder_text(value: object) -> bool:
    text = str(value or "").strip()
    return not text or text == PLACEHOLDER


def _has_source_content(sections: Mapping[str, str]) -> bool:
    for key in ("resumen", "cronologia", "analisis", "riesgos", "recomendaciones", "acciones", "responsables"):
        if not _is_placeholder_text(sections.get(key, PLACEHOLDER)):
            return True
    return False


def _build_prompt(section: str, contexto: str, caso: Mapping[str, object]) -> str:
    categoria = str(caso.get("categoria1") or "Categoría no especificada").strip()
    modalidad = str(caso.get("modalidad") or "Modalidad no especificada").strip()
    canal = str(caso.get("canal") or "Canal no especificado").strip()
    tipo = str(caso.get("tipo_informe") or "Tipo no especificado").strip()
    return (
        "Eres un analista de riesgos que redacta resúmenes ejecutivos en español. "
        "Redacta respuestas de 2 a 4 frases, tono conciso, en voz activa, sin viñetas. "
        "Mantén un máximo de 120 palabras y evita repetir datos ya mencionados. "
        f"Sección objetivo: {section}. "
        f"Tipo de informe: {tipo}; Categoría: {categoria}; Modalidad: {modalidad}; Canal: {canal}. "
        "Usa los datos factuales del caso y la cronología incluida. "
        f"Contexto completo:\n{contexto}"
    )


@dataclass
class SpanishSummaryHelper:
    model_name: str = DEFAULT_MODEL
    max_new_tokens: int = 144
    _pipeline: object = field(default=None, init=False, repr=False)
    _load_error: Exception | None = field(default=None, init=False, repr=False)
    _cache: dict[tuple[str, str, int], str] = field(default_factory=dict, init=False, repr=False)

    def _load_pipeline(self):
        if self._pipeline or self._load_error:
            return
        if not TRANSFORMERS_AVAILABLE:
            self._load_error = RuntimeError("transformers no disponible")
            return
        from transformers import AutoModelForSeq2SeqLM, AutoTokenizer, pipeline

        tokenizer = AutoTokenizer.from_pretrained(self.model_name)
        model = AutoModelForSeq2SeqLM.from_pretrained(self.model_name)
        self._pipeline = pipeline(
            task="text2text-generation",
            model=model,
            tokenizer=tokenizer,
        )

    def summarize(self, section: str, prompt: str, *, max_new_tokens: int | None = None) -> str | None:
        key = (section, prompt, max_new_tokens or self.max_new_tokens)
        if key in self._cache:
            return self._cache[key]

        self._load_pipeline()
        if not self._pipeline:
            return None

        from transformers import set_seed

        try:
            set_seed(0)
            outputs = self._pipeline(
                prompt,
                max_new_tokens=max_new_tokens or self.max_new_tokens,
                num_beams=4,
                do_sample=False,
            )
        except Exception as exc:  # pragma: no cover - defensivo frente a errores del modelo
            self._load_error = exc
            return None

        if not outputs:
            return None
        text = sanitize_rich_text(outputs[0].get("generated_text"), max_chars=MAX_SECTION_CHARS).strip()
        if not text:
            return None
        self._cache[key] = text
        return text


def _synthesize_section_text(
    section: str,
    sections: Mapping[str, str],
    caso: Mapping[str, object],
    llm_helper: SpanishSummaryHelper | None,
) -> str:
    section_key_map = {
        "Resumen": "resumen",
        "Cronología": "cronologia",
        "Cronologia": "cronologia",
        "Análisis": "analisis",
        "Analisis": "analisis",
        "Riesgos": "riesgos",
        "Recomendaciones": "recomendaciones",
        "Recomendación": "recomendaciones",
        "Recomendacion": "recomendaciones",
        "Acciones": "recomendaciones",
        "Responsables": "responsables",
    }
    section_key = section_key_map.get(section, section.lower())
    fallback_source = sections.get(section_key, PLACEHOLDER)
    if _is_placeholder_text(fallback_source):
        return PLACEHOLDER
    if not _has_source_content(sections):
        return PLACEHOLDER
    if not llm_helper:
        return fallback_source
    context_lines = [
        f"Caso: {sections.get('codigo', PLACEHOLDER)} | Tipo: {caso.get('tipo_informe', PLACEHOLDER)}",
        f"Resumen: {sections.get('resumen', PLACEHOLDER)}",
        f"Cronología: {sections.get('cronologia', PLACEHOLDER)}",
        f"Análisis: {sections.get('analisis', PLACEHOLDER)}",
        f"Riesgos: {sections.get('riesgos', PLACEHOLDER)}",
        f"Recomendaciones: {sections.get('recomendaciones', sections.get('acciones', PLACEHOLDER))}",
        f"Responsables: {sections.get('responsables', PLACEHOLDER)}",
    ]
    prompt = _build_prompt(section, "\n".join(context_lines), caso)
    llm_summary = llm_helper.summarize(section, prompt)
    return llm_summary or fallback_source


def _add_section_panel(slide, left, top, width, height, title: str, body: str, *, accent: RGBColor | None = None):
    if not PPTX_AVAILABLE:
        return None
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(248, 250, 253)
    panel.line.color.rgb = RGBColor(196, 204, 219)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, SECTION_HEADER_HEIGHT)
    header.fill.solid()
    header.fill.fore_color.rgb = accent or RGBColor(230, 233, 240)
    header.line.fill.background()
    header_frame = header.text_frame
    header_frame.clear()
    header_para = header_frame.paragraphs[0]
    header_para.text = title
    header_para.font.size = Pt(11)
    header_para.font.bold = True
    header_para.font.color.rgb = RGBColor(32, 45, 71)

    text_left = left + PANEL_PADDING
    text_top = top + SECTION_HEADER_HEIGHT + PANEL_PADDING
    text_width = width - (2 * PANEL_PADDING)
    text_height = height - SECTION_HEADER_HEIGHT - (2 * PANEL_PADDING)
    fit = _fit_text_to_box(
        body,
        _emu_to_inches(text_width),
        _emu_to_inches(text_height),
        section_title=title,
    )
    body_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    frame = body_box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraphs = _split_body_paragraphs(fit.text)
    for index, (text, is_bullet) in enumerate(paragraphs):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.text = text
        para.font.size = Pt(fit.font_pt)
        para.alignment = PP_ALIGN.LEFT
        para.line_spacing = DEFAULT_LINE_SPACING
        _set_paragraph_bullet(para, is_bullet)
    return panel


def _add_masthead(
    slide,
    slide_width,
    title: str,
    case_text: str,
    codigo_text: str,
    issuer_text: str,
):
    if not PPTX_AVAILABLE:
        return None
    masthead = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, MASTHEAD_HEIGHT)
    masthead.fill.solid()
    masthead.fill.fore_color.rgb = RGBColor(21, 43, 77)
    masthead.line.fill.background()

    left_box = slide.shapes.add_textbox(MARGIN, Inches(0.12), int(slide_width * 0.62), MASTHEAD_HEIGHT - Inches(0.2))
    left_frame = left_box.text_frame
    left_frame.clear()
    title_para = left_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(11)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)

    case_line = left_frame.add_paragraph()
    case_line.text = f"Caso: {case_text}"
    case_line.font.size = Pt(11)
    case_line.font.color.rgb = RGBColor(214, 225, 240)

    right_box_width = int(slide_width * 0.32)
    right_box = slide.shapes.add_textbox(slide_width - right_box_width - MARGIN, Inches(0.12), right_box_width, MASTHEAD_HEIGHT - Inches(0.24))
    right_frame = right_box.text_frame
    right_frame.clear()
    code_para = right_frame.paragraphs[0]
    code_para.text = f"Código: {codigo_text}"
    code_para.font.size = Pt(11)
    code_para.font.bold = True
    code_para.font.color.rgb = RGBColor(255, 255, 255)
    code_para.alignment = PP_ALIGN.RIGHT

    issuer_para = right_frame.add_paragraph()
    issuer_para.text = f"Emitido por: {issuer_text}"
    issuer_para.font.size = Pt(11)
    issuer_para.font.color.rgb = RGBColor(214, 225, 240)
    issuer_para.alignment = PP_ALIGN.RIGHT
    return masthead


def _add_executive_summary_slide(slide, sections: Mapping[str, str], summary: ExecutiveSummary):
    if not PPTX_AVAILABLE:
        return None
    _add_masthead(
        slide,
        SLIDE_WIDTH_16_9,
        "Resumen ejecutivo · Alerta temprana",
        sections.get("caso", PLACEHOLDER),
        sections.get("codigo", PLACEHOLDER),
        sections.get("emitido_por", PLACEHOLDER),
    )
    body_top = MARGIN + MASTHEAD_HEIGHT + Inches(0.05)
    full_width = SLIDE_WIDTH_16_9 - (2 * MARGIN)
    available_height = SLIDE_HEIGHT_16_9 - body_top - MARGIN
    gap = Inches(0.08)
    panel_height = (available_height - (2 * gap)) / 3

    _add_section_panel(
        slide,
        MARGIN,
        body_top,
        full_width,
        panel_height,
        "Mensaje clave",
        summary.headline,
        accent=RGBColor(219, 223, 232),
    )
    _add_section_panel(
        slide,
        MARGIN,
        body_top + panel_height + gap,
        full_width,
        panel_height,
        "Puntos de soporte (3-5)",
        "\n".join(f"• {line}" for line in summary.supporting_points),
        accent=RGBColor(219, 223, 232),
    )
    _add_section_panel(
        slide,
        MARGIN,
        body_top + (2 * panel_height) + (2 * gap),
        full_width,
        panel_height,
        "Evidencia / trazabilidad",
        "\n".join(f"• {line}" for line in summary.evidence),
        accent=RGBColor(214, 226, 240),
    )


def build_alerta_temprana_ppt(
    data: CaseData | Mapping[str, object],
    output_path: Path,
    llm_helper: SpanishSummaryHelper | None = None,
) -> Path:
    """Genera una presentación PPTX con la alerta temprana."""

    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx no está disponible para generar la alerta temprana.")

    dataset = data if isinstance(data, CaseData) else CaseData.from_mapping(data or {})
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    caso = dataset.get("caso", {}) if isinstance(dataset, Mapping) else {}
    sections = build_alerta_temprana_sections(dataset)
    resumen_ejecutivo = build_executive_summary(dataset)
    llm = llm_helper

    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH_16_9
    presentation.slide_height = SLIDE_HEIGHT_16_9
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_executive_summary_slide(slide, sections, resumen_ejecutivo)

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_masthead(
        slide,
        presentation.slide_width,
        sections.get("titulo_reporte", "Reporte de Alertas Tempranas por Casos de Fraude"),
        sections.get("caso", PLACEHOLDER),
        sections.get("codigo", PLACEHOLDER),
        sections.get("emitido_por", PLACEHOLDER),
    )

    body_top = MARGIN + MASTHEAD_HEIGHT + Inches(0.05)
    total_width = presentation.slide_width - (2 * MARGIN) - COLUMN_GAP
    left_column_width = int(total_width * 0.66)
    right_column_width = total_width - left_column_width

    left_x = MARGIN
    right_x = MARGIN + left_column_width + COLUMN_GAP

    available_height = presentation.slide_height - body_top - MARGIN
    section_gap = Inches(0.08)

    left_usable_height = available_height - (2 * section_gap)
    resumen_height = int(left_usable_height * 0.36)
    cronologia_height = int(left_usable_height * 0.3)
    analisis_height = left_usable_height - resumen_height - cronologia_height

    right_usable_height = available_height - (2 * section_gap)
    riesgos_height = int(right_usable_height * 0.34)
    acciones_height = int(right_usable_height * 0.26)
    responsables_height = max(right_usable_height - riesgos_height - acciones_height, int(Inches(1)))

    resumen_text = _synthesize_section_text("Resumen", sections, caso, llm)
    _add_section_panel(
        slide,
        left_x,
        body_top,
        left_column_width,
        resumen_height,
        "Resumen",
        resumen_text,
        accent=RGBColor(219, 223, 232),
    )

    cronologia_text = _synthesize_section_text("Cronología", sections, caso, llm)
    _add_section_panel(
        slide,
        left_x,
        body_top + resumen_height + section_gap,
        left_column_width,
        cronologia_height,
        "Cronología",
        cronologia_text,
        accent=RGBColor(219, 223, 232),
    )

    analisis_text = _synthesize_section_text("Análisis", sections, caso, llm)
    _add_section_panel(
        slide,
        left_x,
        body_top + resumen_height + cronologia_height + (2 * section_gap),
        left_column_width,
        analisis_height,
        "Análisis",
        analisis_text,
        accent=RGBColor(219, 223, 232),
    )

    riesgos_text = _synthesize_section_text("Riesgos", sections, caso, llm)
    _add_section_panel(
        slide,
        right_x,
        body_top,
        right_column_width,
        riesgos_height,
        "Riesgos identificados",
        riesgos_text,
        accent=RGBColor(214, 226, 240),
    )

    recomendaciones_text = _synthesize_section_text("Recomendaciones", sections, caso, llm)
    _add_section_panel(
        slide,
        right_x,
        body_top + riesgos_height + section_gap,
        right_column_width,
        acciones_height,
        "Recomendaciones",
        recomendaciones_text,
        accent=RGBColor(214, 226, 240),
    )

    responsables_text = _synthesize_section_text("Responsables", sections, caso, llm)
    _add_section_panel(
        slide,
        right_x,
        body_top + riesgos_height + acciones_height + (2 * section_gap),
        right_column_width,
        responsables_height,
        "Responsables",
        responsables_text,
        accent=RGBColor(214, 226, 240),
    )

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "Plantilla 16:9 con barra superior azul oscuro que expone el tipo de reporte, el caso y los campos de código/"
        "emisor. El cuerpo usa grilla de dos columnas: izquierda (2/3) con Resumen, Cronología y Análisis apilados;"
        " derecha (1/3) con Riesgos identificados, Recomendaciones y Responsables. Las barras grises de sección"
        " indican los encabezados. Se prioriza texto determinístico: resumen con montos, cronología desde hallazgos,"
        " análisis desde antecedentes/conclusiones y riesgos desde catálogo."
    )

    presentation.save(output_path)
    return output_path


__all__ = [
    "build_alerta_temprana_ppt",
    "SpanishSummaryHelper",
    "PPTX_AVAILABLE",
    "PPTX_MISSING_MESSAGE",
]
