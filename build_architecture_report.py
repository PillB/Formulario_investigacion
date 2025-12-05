"""Utility to render the architecture PDF for Formulario_investigacion.

The script regenerates Mermaid diagrams (PNG) using mermaid-cli and then
assembles a styled PDF with ReportLab. It keeps all logic in one place so the
repo does not need to track the binary PDF artifact.
"""
from __future__ import annotations

import argparse
import datetime as _dt
import shutil
import subprocess
import sys
import textwrap
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A3, LETTER, landscape
from reportlab.lib.styles import ParagraphStyle, StyleSheet1, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.lib.utils import ImageReader
from reportlab.platypus import (
    BaseDocTemplate,
    Frame,
    Image,
    ListFlowable,
    ListItem,
    NextPageTemplate,
    PageBreak,
    PageTemplate,
    Paragraph,
    Spacer,
    Table,
)
from reportlab.platypus.tableofcontents import TableOfContents

PROJECT_ROOT = Path(__file__).parent
DOCS_DIR = PROJECT_ROOT / "docs"
DEFAULT_OUTPUT = PROJECT_ROOT / "Formulario_Investigacion_Architecture_and_Data_Flow.pdf"
DEFAULT_PPTX = PROJECT_ROOT / "Formulario_Investigacion_Diagramas_editables.pptx"
ARCH_MMD = DOCS_DIR / "architecture.mmd"
ARCH_PNG = DOCS_DIR / "architecture.png"
DB_ARCH_MMD = DOCS_DIR / "db_architecture.mmd"
DB_ARCH_PNG = DOCS_DIR / "db_architecture.png"
SEQ_MMD = DOCS_DIR / "sequence_diagram.mmd"
SEQ_PNG = DOCS_DIR / "sequence.png"
PUPPETEER_CONFIG = PROJECT_ROOT / "puppeteer-config.json"
MERMAID_CONFIG = PROJECT_ROOT / "mermaid-config.json"
MERMAID_EXPORT_SCALE = 2.8  # higher pixel density to keep diagram text crisp
MERMAID_EXPORT_WIDTH_PX = 4800  # approximate A3 width at ~290 DPI

PPTX_SLIDE_WIDTH = Cm(42)  # A3 landscape width
PPTX_SLIDE_HEIGHT = Cm(29.7)  # A3 landscape height
PPTX_MARGIN = Cm(1.2)
PPTX_TITLE_HEIGHT = Cm(2.2)
PPTX_SUBTITLE_HEIGHT = Cm(1.1)
PPTX_CONTENT_GAP = Cm(0.5)
PPTX_BODY_TEXT_SIZE = Pt(18)
PPTX_TITLE_TEXT_SIZE = Pt(32)
PPTX_SUBTITLE_TEXT_SIZE = Pt(18)
PPTX_HEADER_FILL = RGBColor(237, 242, 247)
PPTX_CONTEXT_FILL = RGBColor(225, 242, 255)
PPTX_APP_FILL = RGBColor(236, 249, 235)
PPTX_COMPONENT_FILL = RGBColor(252, 243, 229)
PPTX_SERVICE_FILL = RGBColor(231, 231, 255)
PPTX_BORDER_COLOR = RGBColor(64, 64, 64)
PPTX_TEXT_COLOR = RGBColor(0, 0, 0)


class HeadingParagraph(Paragraph):
    """Paragraph that stores the heading level to feed the TOC."""

    def __init__(self, text: str, style: ParagraphStyle, level: int):
        super().__init__(text, style)
        self.toc_level = level


# ---------------------------------------------------------------------------
# Mermaid helpers
# ---------------------------------------------------------------------------


def render_mermaid(source: Path, target: Path) -> Path:
    """Render a Mermaid file to PNG using mermaid-cli.

    The function tries an installed ``mmdc`` binary first and falls back to
    ``npx -y @mermaid-js/mermaid-cli`` to avoid tracking rendered assets in git.
    """

    if not source.exists():
        raise FileNotFoundError(source)

    target.parent.mkdir(parents=True, exist_ok=True)

    cmd: list[str]
    if shutil.which("mmdc"):
        cmd = ["mmdc", "-i", str(source), "-o", str(target)]
    elif shutil.which("npx"):
        cmd = [
            "npx",
            "-y",
            "@mermaid-js/mermaid-cli",
            "-i",
            str(source),
            "-o",
            str(target),
        ]
    else:
        raise RuntimeError(
            "Se requiere mermaid-cli. Instala 'npm install -g @mermaid-js/mermaid-cli'"
        )

    cmd.extend(["-s", str(MERMAID_EXPORT_SCALE), "-w", str(MERMAID_EXPORT_WIDTH_PX)])
    if PUPPETEER_CONFIG.exists():
        cmd.extend(["-p", str(PUPPETEER_CONFIG)])
    if MERMAID_CONFIG.exists():
        cmd.extend(["-c", str(MERMAID_CONFIG)])

    subprocess.run(cmd, check=True)
    return target


def _ensure_db_diagram_png() -> Path:
    """Render the DB Mermaid diagram or create a placeholder if CLI is missing."""

    try:
        return render_mermaid(DB_ARCH_MMD, DB_ARCH_PNG)
    except RuntimeError:
        try:
            from PIL import Image, ImageDraw, ImageFont  # type: ignore

            DB_ARCH_PNG.parent.mkdir(parents=True, exist_ok=True)
            img = Image.new("RGB", (1200, 800), color="white")
            draw = ImageDraw.Draw(img)
            draw.rectangle((20, 20, 1180, 780), outline="black", width=4)
            draw.text(
                (60, 60),
                "DB diagram placeholder\nInstala mermaid-cli para render completo",
                fill="black",
                font=ImageFont.load_default(),
            )
            img.save(DB_ARCH_PNG)
            return DB_ARCH_PNG
        except Exception as exc:  # pragma: no cover - fallback only
            raise RuntimeError(
                "Se requiere mermaid-cli. Instala 'npm install -g @mermaid-js/mermaid-cli'"
            ) from exc


# ---------------------------------------------------------------------------
# Styling
# ---------------------------------------------------------------------------


def _build_stylesheet() -> StyleSheet1:
    styles = getSampleStyleSheet()

    def _ensure_style(name: str, **attributes: object) -> ParagraphStyle:
        style = styles[name] if name in styles.byName else ParagraphStyle(name=name)
        for attr, value in attributes.items():
            setattr(style, attr, value)
        if name not in styles.byName:
            styles.add(style)
        return style

    _ensure_style(
        name="CoverTitle",
        fontSize=24,
        leading=28,
        spaceAfter=18,
        alignment=1,
    )
    _ensure_style(
        name="CoverSubtitle",
        fontSize=12,
        leading=14,
        textColor=colors.grey,
        alignment=1,
        spaceAfter=6,
    )
    _ensure_style(
        name="Meta",
        fontSize=10,
        leading=12,
        textColor=colors.HexColor("#555555"),
        alignment=1,
        spaceAfter=14,
    )

    _ensure_style(name="Heading1", fontSize=18, leading=22, spaceAfter=12)
    _ensure_style(name="Heading2", fontSize=14, leading=18, spaceAfter=8)

    _ensure_style(
        name="Body",
        parent=styles["BodyText"],
        fontSize=10.5,
        leading=14,
        spaceAfter=8,
    )
    _ensure_style(
        name="Bullet",
        parent=styles["Body"],
        leftIndent=12,
        bulletIndent=0,
        spaceAfter=4,
    )
    _ensure_style(
        name="TableHeader",
        parent=styles["Body"],
        fontSize=10,
        textColor=colors.white,
        backColor=colors.HexColor("#1f3a93"),
        spaceAfter=4,
    )
    _ensure_style(
        name="TableCell",
        parent=styles["Body"],
        fontSize=9,
        leading=12,
    )
    return styles


# ---------------------------------------------------------------------------
# PDF assembly
# ---------------------------------------------------------------------------


def _after_flowable(doc: BaseDocTemplate, flowable):
    if getattr(flowable, "toc_level", None) is not None:
        text = flowable.getPlainText()
        level = flowable.toc_level
        doc.notify("TOCEntry", (level, text, doc.page))


def _page_decor(canvas, doc: BaseDocTemplate):
    canvas.saveState()
    canvas.setFont("Helvetica", 9)
    canvas.setFillColor(colors.HexColor("#666666"))
    canvas.drawRightString(
        doc.pagesize[0] - doc.rightMargin,
        doc.bottomMargin - 12,
        f"Página {doc.page}",
    )
    canvas.restoreState()


def _heading(text: str, styles: StyleSheet1, level: int) -> HeadingParagraph:
    style_name = f"Heading{level}"
    return HeadingParagraph(text, styles[style_name], level)


def _paragraph(text: str, styles: StyleSheet1) -> Paragraph:
    return Paragraph(textwrap.dedent(text).strip(), styles["Body"])


def _bullet_list(items: list[str], styles: StyleSheet1) -> ListFlowable:
    return ListFlowable(
        [ListItem(_paragraph(item, styles), leftIndent=0) for item in items],
        bulletType="bullet",
        start=None,
        leftIndent=12,
    )


def _technology_table(styles: StyleSheet1) -> Table:
    data = [
        ["Capa", "Tecnología"],
        ["UI de escritorio", "Tkinter / ttk"],
        ["Diagramas", "Mermaid (architecture.mmd, sequence_diagram.mmd)"],
        ["Persistencia temporal", "autosave.json, autosaves/<caso>/auto_<n>.json"],
        ["Validaciones", "validators.py y FieldValidator"],
        ["Importación", "utils.iter_massive_csv_rows y importadores en Acciones"],
        ["Reportes", "report_builder.py (CSV/JSON/Markdown/DOCX)"],
        ["Logs y analítica", "logs.csv + analytics/usage_visualizer.py"],
    ]

    table = Table(data, repeatRows=1)
    table.setStyle(
        [
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1f3a93")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("ALIGN", (0, 0), (-1, 0), "CENTER"),
            ("ALIGN", (0, 1), (-1, -1), "LEFT"),
            ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("BACKGROUND", (0, 1), (-1, -1), colors.whitesmoke),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.whitesmoke, colors.white]),
            ("GRID", (0, 0), (-1, -1), 0.25, colors.lightgrey),
            ("LEFTPADDING", (0, 0), (-1, -1), 6),
            ("RIGHTPADDING", (0, 0), (-1, -1), 6),
            ("TOPPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ]
    )
    return table


def _flowable_image(image_path: Path, target_width: float) -> Image:
    """Scale an image to the requested width while keeping the aspect ratio."""

    reader = ImageReader(str(image_path))
    original_width, original_height = reader.getSize()
    if original_width <= 0 or original_height <= 0:
        raise ValueError("Image dimensions must be positive")

    scale = target_width / float(original_width)
    target_height = original_height * scale
    return Image(str(image_path), width=target_width, height=target_height)


def _scale_image_to_box(image_path: Path, max_width: float, max_height: float) -> tuple[float, float]:
    """Return width/height scaled to fit within a bounding box."""

    reader = ImageReader(str(image_path))
    original_width, original_height = reader.getSize()
    if original_width <= 0 or original_height <= 0:
        raise ValueError("Image dimensions must be positive")

    width_scale = max_width / float(original_width)
    height_scale = max_height / float(original_height)
    scale = min(width_scale, height_scale)
    return original_width * scale, original_height * scale


# ---------------------------------------------------------------------------
# PowerPoint helpers (native shapes)
# ---------------------------------------------------------------------------


def _add_header(slide, *, title: str, subtitle: str, slide_width: float) -> float:
    title_box = slide.shapes.add_textbox(
        PPTX_MARGIN,
        PPTX_MARGIN,
        slide_width - 2 * PPTX_MARGIN,
        PPTX_TITLE_HEIGHT,
    )
    title_box.name = "header_title"
    title_tf = title_box.text_frame
    title_tf.paragraphs[0].text = title
    title_tf.paragraphs[0].font.size = PPTX_TITLE_TEXT_SIZE
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = PPTX_TEXT_COLOR
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(
        PPTX_MARGIN,
        PPTX_MARGIN + PPTX_TITLE_HEIGHT,
        slide_width - 2 * PPTX_MARGIN,
        PPTX_SUBTITLE_HEIGHT,
    )
    subtitle_box.name = "header_subtitle"
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.word_wrap = True
    subtitle_tf.paragraphs[0].text = subtitle
    subtitle_tf.paragraphs[0].font.size = PPTX_SUBTITLE_TEXT_SIZE
    subtitle_tf.paragraphs[0].font.color.rgb = PPTX_TEXT_COLOR
    subtitle_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return PPTX_MARGIN + PPTX_TITLE_HEIGHT + PPTX_SUBTITLE_HEIGHT + PPTX_CONTENT_GAP


def _add_box(
    slide,
    *,
    name: str,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    fill: RGBColor,
    font_size: Pt = PPTX_BODY_TEXT_SIZE,
    bold: bool = False,
    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
    align: PP_ALIGN = PP_ALIGN.CENTER,
):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = PPTX_BORDER_COLOR
    shape.text_frame.word_wrap = True
    para = shape.text_frame.paragraphs[0]
    para.text = text
    para.font.size = font_size
    para.font.bold = bold
    para.font.color.rgb = PPTX_TEXT_COLOR
    para.alignment = align
    return shape


def _anchor_point(shape, anchor: str) -> tuple[float, float]:
    match anchor:
        case "top":
            return shape.left + shape.width / 2, shape.top
        case "bottom":
            return shape.left + shape.width / 2, shape.top + shape.height
        case "left":
            return shape.left, shape.top + shape.height / 2
        case "right":
            return shape.left + shape.width, shape.top + shape.height / 2
        case _:
            return shape.left + shape.width / 2, shape.top + shape.height / 2


def _connect(
    shapes,
    start_shape,
    end_shape,
    *,
    name: str,
    start_anchor: str = "center",
    end_anchor: str = "center",
    connector_type: MSO_CONNECTOR = MSO_CONNECTOR.ELBOW,
):
    start_x, start_y = _anchor_point(start_shape, start_anchor)
    end_x, end_y = _anchor_point(end_shape, end_anchor)
    connector = shapes.add_connector(connector_type, start_x, start_y, end_x, end_y)
    connector.name = name
    connector.line.width = Pt(2)
    return connector


def _add_note(slide, *, text: str, left: float, top: float, width: float, height: float):
    note = _add_box(
        slide,
        name=f"note_{text[:10]}",
        text=text,
        left=left,
        top=top,
        width=width,
        height=height,
        fill=PPTX_HEADER_FILL,
        font_size=Pt(16),
        shape_type=MSO_SHAPE.RECTANGLE,
        align=PP_ALIGN.LEFT,
    )
    note.line.color.rgb = PPTX_BORDER_COLOR
    note.line.dash_style = MSO_LINE_DASH_STYLE.DASH_DOT
    return note


def _architecture_nodes(slide, *, content_top: float, slide_width: float) -> dict[str, object]:
    """Create the architecture diagram with native shapes on a 100x100 grid."""

    unit_x = (slide_width - 2 * PPTX_MARGIN) / 100
    unit_y = (PPTX_SLIDE_HEIGHT - content_top - PPTX_MARGIN) / 100

    def g(x_units: float, y_units: float, w_units: float, h_units: float) -> tuple[float, float, float, float]:
        return (
            PPTX_MARGIN + x_units * unit_x,
            content_top + y_units * unit_y,
            w_units * unit_x,
            h_units * unit_y,
        )

    nodes: dict[str, object] = {}

    # Contexto en cuadrícula 100x100
    nodes["gestion"] = _add_box(
        slide,
        name="context_gestion",
        text="Gestión de Investigaciones",
        left=g(0, 0, 14, 8)[0],
        top=g(0, 0, 14, 8)[1],
        width=g(0, 0, 14, 8)[2],
        height=g(0, 0, 14, 8)[3],
        fill=PPTX_CONTEXT_FILL,
        bold=True,
    )
    nodes["analista"] = _add_box(
        slide,
        name="context_analista",
        text="Investigador",
        left=g(16, 0, 14, 8)[0],
        top=g(16, 0, 14, 8)[1],
        width=g(16, 0, 14, 8)[2],
        height=g(16, 0, 14, 8)[3],
        fill=PPTX_CONTEXT_FILL,
    )
    nodes["tkapp"] = _add_box(
        slide,
        name="context_tkapp",
        text="Aplicación Python Tkinter GUI\\nFormularioInvestigacionesFraude",
        left=g(34, 0, 20, 10)[0],
        top=g(34, 0, 20, 10)[1],
        width=g(34, 0, 20, 10)[2],
        height=g(34, 0, 20, 10)[3],
        fill=PPTX_CONTEXT_FILL,
        bold=True,
    )
    nodes["databricks"] = _add_box(
        slide,
        name="context_databricks",
        text="Catálogos empresariales futuro\\nDatabricks Parquet/CSV",
        left=g(58, 0, 16, 9)[0],
        top=g(58, 0, 16, 9)[1],
        width=g(58, 0, 16, 9)[2],
        height=g(58, 0, 16, 9)[3],
        fill=PPTX_CONTEXT_FILL,
    )
    nodes["grc"] = _add_box(
        slide,
        name="context_grc",
        text="Registro de riesgos GRC",
        left=g(78, 0, 10, 8)[0],
        top=g(78, 0, 10, 8)[1],
        width=g(78, 0, 10, 8)[2],
        height=g(78, 0, 10, 8)[3],
        fill=PPTX_CONTEXT_FILL,
    )
    nodes["normdb"] = _add_box(
        slide,
        name="context_normdb",
        text="Índice maestro de normas",
        left=g(90, 0, 10, 8)[0],
        top=g(90, 0, 10, 8)[1],
        width=g(90, 0, 10, 8)[2],
        height=g(90, 0, 10, 8)[3],
        fill=PPTX_CONTEXT_FILL,
    )

    # Contenedores principales
    nodes["main"] = _add_box(
        slide,
        name="container_main",
        text="main.py\\nArranque Tk/Ttk",
        left=g(10, 16, 16, 9)[0],
        top=g(10, 16, 16, 9)[1],
        width=g(10, 16, 16, 9)[2],
        height=g(10, 16, 16, 9)[3],
        fill=PPTX_APP_FILL,
    )
    nodes["appcore"] = _add_box(
        slide,
        name="container_appcore",
        text="FormularioInvestigacionesFraude app.py\\nControlador",
        left=g(30, 16, 18, 11)[0],
        top=g(30, 16, 18, 11)[1],
        width=g(30, 16, 18, 11)[2],
        height=g(30, 16, 18, 11)[3],
        fill=PPTX_APP_FILL,
        bold=True,
    )
    nodes["ui"] = _add_box(
        slide,
        name="container_ui",
        text="ui/frames/*\\nPestañas de Caso, Clientes, Team Members, Productos, Riesgos, Normas",
        left=g(52, 16, 20, 13)[0],
        top=g(52, 16, 20, 13)[1],
        width=g(52, 16, 20, 13)[2],
        height=g(52, 16, 20, 13)[3],
        fill=PPTX_APP_FILL,
    )

    nodes["services"] = _add_box(
        slide,
        name="container_services",
        text="models/*\\nCatálogos y autopoblado",
        left=g(0, 32, 18, 10)[0],
        top=g(0, 32, 18, 10)[1],
        width=g(0, 32, 18, 10)[2],
        height=g(0, 32, 18, 10)[3],
        fill=PPTX_SERVICE_FILL,
    )
    nodes["validators"] = _add_box(
        slide,
        name="container_validators",
        text="validators.py\\nReglas + logs",
        left=g(20, 32, 16, 10)[0],
        top=g(20, 32, 16, 10)[1],
        width=g(20, 32, 16, 10)[2],
        height=g(20, 32, 16, 10)[3],
        fill=PPTX_SERVICE_FILL,
    )
    nodes["reporting"] = _add_box(
        slide,
        name="container_reporting",
        text="report_builder.py\\nGeneración JSON/CSV/MD/DOCX",
        left=g(38, 32, 18, 10)[0],
        top=g(38, 32, 18, 10)[1],
        width=g(38, 32, 18, 10)[2],
        height=g(38, 32, 18, 10)[3],
        fill=PPTX_SERVICE_FILL,
    )
    nodes["files"] = _add_box(
        slide,
        name="container_files",
        text="Sistema de archivos\\nautosave.json, autosaves/<caso>/auto_N.json, exports/",
        left=g(58, 32, 22, 11)[0],
        top=g(58, 32, 22, 11)[1],
        width=g(58, 32, 22, 11)[2],
        height=g(58, 32, 22, 11)[3],
        fill=PPTX_SERVICE_FILL,
    )
    nodes["external"] = _add_box(
        slide,
        name="container_external",
        text="Unidad externa\\nbackups y logs",
        left=g(82, 32, 14, 9)[0],
        top=g(82, 32, 14, 9)[1],
        width=g(82, 32, 14, 9)[2],
        height=g(82, 32, 14, 9)[3],
        fill=PPTX_SERVICE_FILL,
    )

    # Componentes de aplicación
    nodes["caseTab"] = _add_box(
        slide,
        name="component_case",
        text="Pestaña Caso",
        left=g(2, 48, 14, 8)[0],
        top=g(2, 48, 14, 8)[1],
        width=g(2, 48, 14, 8)[2],
        height=g(2, 48, 14, 8)[3],
        fill=PPTX_COMPONENT_FILL,
    )
    nodes["clientTab"] = _add_box(
        slide,
        name="component_client",
        text="Pestaña Clientes",
        left=g(18, 48, 14, 8)[0],
        top=g(18, 48, 14, 8)[1],
        width=g(18, 48, 14, 8)[2],
        height=g(18, 48, 14, 8)[3],
        fill=PPTX_COMPONENT_FILL,
    )
    nodes["teamTab"] = _add_box(
        slide,
        name="component_team",
        text="Pestaña Team Members",
        left=g(34, 48, 14, 8)[0],
        top=g(34, 48, 14, 8)[1],
        width=g(34, 48, 14, 8)[2],
        height=g(34, 48, 14, 8)[3],
        fill=PPTX_COMPONENT_FILL,
    )
    nodes["productTab"] = _add_box(
        slide,
        name="component_product",
        text="Pestaña Productos",
        left=g(50, 48, 14, 8)[0],
        top=g(50, 48, 14, 8)[1],
        width=g(50, 48, 14, 8)[2],
        height=g(50, 48, 14, 8)[3],
        fill=PPTX_COMPONENT_FILL,
    )
    nodes["riskTab"] = _add_box(
        slide,
        name="component_risk",
        text="Pestaña Riesgos",
        left=g(66, 48, 14, 8)[0],
        top=g(66, 48, 14, 8)[1],
        width=g(66, 48, 14, 8)[2],
        height=g(66, 48, 14, 8)[3],
        fill=PPTX_COMPONENT_FILL,
    )

    nodes["normTab"] = _add_box(
        slide,
        name="component_norm",
        text="Pestaña Normas",
        left=g(82, 48, 14, 8)[0],
        top=g(82, 48, 14, 8)[1],
        width=g(82, 48, 14, 8)[2],
        height=g(82, 48, 14, 8)[3],
        fill=PPTX_COMPONENT_FILL,
    )
    nodes["actions"] = _add_box(
        slide,
        name="component_actions",
        text="Tab Acciones",
        left=g(22, 60, 14, 8)[0],
        top=g(22, 60, 14, 8)[1],
        width=g(22, 60, 14, 8)[2],
        height=g(22, 60, 14, 8)[3],
        fill=PPTX_COMPONENT_FILL,
    )
    nodes["summary"] = _add_box(
        slide,
        name="component_summary",
        text="Tab Resumen",
        left=g(42, 60, 14, 8)[0],
        top=g(42, 60, 14, 8)[1],
        width=g(42, 60, 14, 8)[2],
        height=g(42, 60, 14, 8)[3],
        fill=PPTX_COMPONENT_FILL,
    )

    nodes["autosave"] = _add_box(
        slide,
        name="component_autosave",
        text="Auto guardado temporal",
        left=g(62, 60, 14, 8)[0],
        top=g(62, 60, 14, 8)[1],
        width=g(62, 60, 14, 8)[2],
        height=g(62, 60, 14, 8)[3],
        fill=PPTX_COMPONENT_FILL,
    )
    nodes["logging"] = _add_box(
        slide,
        name="component_logging",
        text="Bitácora logs.csv",
        left=g(78, 60, 14, 8)[0],
        top=g(78, 60, 14, 8)[1],
        width=g(78, 60, 14, 8)[2],
        height=g(78, 60, 14, 8)[3],
        fill=PPTX_COMPONENT_FILL,
    )
    nodes["importer"] = _add_box(
        slide,
        name="component_importer",
        text="Importadores CSV",
        left=g(6, 72, 16, 9)[0],
        top=g(6, 72, 16, 9)[1],
        width=g(6, 72, 16, 9)[2],
        height=g(6, 72, 16, 9)[3],
        fill=PPTX_COMPONENT_FILL,
    )
    nodes["autofill"] = _add_box(
        slide,
        name="component_autofill",
        text="Servicio de Catálogos\\nAutopoblado",
        left=g(26, 72, 18, 9)[0],
        top=g(26, 72, 18, 9)[1],
        width=g(26, 72, 18, 9)[2],
        height=g(26, 72, 18, 9)[3],
        fill=PPTX_COMPONENT_FILL,
    )
    nodes["reporter"] = _add_box(
        slide,
        name="component_reporter",
        text="report_builder",
        left=g(48, 72, 16, 9)[0],
        top=g(48, 72, 16, 9)[1],
        width=g(48, 72, 16, 9)[2],
        height=g(48, 72, 16, 9)[3],
        fill=PPTX_COMPONENT_FILL,
    )

    # Relacionar
    _connect(
        slide.shapes,
        nodes["gestion"],
        nodes["analista"],
        name="gestion_to_analista",
        start_anchor="right",
        end_anchor="left",
    )
    _connect(
        slide.shapes,
        nodes["analista"],
        nodes["tkapp"],
        name="analista_to_tkapp",
        start_anchor="right",
        end_anchor="left",
    )
    _connect(
        slide.shapes,
        nodes["databricks"],
        nodes["tkapp"],
        name="databricks_to_tkapp",
        start_anchor="bottom",
        end_anchor="right",
    )
    _connect(
        slide.shapes,
        nodes["grc"],
        nodes["tkapp"],
        name="grc_to_tkapp",
        start_anchor="bottom",
        end_anchor="right",
    )
    _connect(
        slide.shapes,
        nodes["normdb"],
        nodes["tkapp"],
        name="normdb_to_tkapp",
        start_anchor="top",
        end_anchor="right",
    )

    _connect(
        slide.shapes,
        nodes["tkapp"],
        nodes["main"],
        name="tkapp_to_main",
        start_anchor="bottom",
        end_anchor="top",
    )
    _connect(
        slide.shapes,
        nodes["main"],
        nodes["appcore"],
        name="main_to_appcore",
        start_anchor="right",
        end_anchor="left",
    )
    _connect(
        slide.shapes,
        nodes["appcore"],
        nodes["ui"],
        name="appcore_to_ui",
        start_anchor="right",
        end_anchor="left",
    )
    _connect(
        slide.shapes,
        nodes["appcore"],
        nodes["services"],
        name="appcore_to_services",
        start_anchor="left",
        end_anchor="top",
    )
    _connect(
        slide.shapes,
        nodes["appcore"],
        nodes["validators"],
        name="appcore_to_validators",
        start_anchor="left",
        end_anchor="top",
    )
    _connect(
        slide.shapes,
        nodes["appcore"],
        nodes["reporting"],
        name="appcore_to_reporting",
        start_anchor="bottom",
        end_anchor="top",
    )
    _connect(
        slide.shapes,
        nodes["reporting"],
        nodes["files"],
        name="reporting_to_files",
        start_anchor="right",
        end_anchor="left",
    )
    _connect(
        slide.shapes,
        nodes["files"],
        nodes["external"],
        name="files_to_external",
        start_anchor="right",
        end_anchor="left",
    )

    # Pestañas
    for tab in [
        "caseTab",
        "clientTab",
        "teamTab",
        "productTab",
        "riskTab",
        "normTab",
        "actions",
        "summary",
    ]:
        _connect(
            slide.shapes,
            nodes["appcore"],
            nodes[tab],
            name=f"appcore_to_{tab}",
            start_anchor="bottom",
            end_anchor="top",
        )
        _connect(
            slide.shapes,
            nodes["importer"],
            nodes[tab],
            name=f"importer_to_{tab}",
            start_anchor="top",
            end_anchor="bottom",
        )
        _connect(
            slide.shapes,
            nodes[tab],
            nodes["autofill"],
            name=f"{tab}_to_autofill",
            start_anchor="right",
            end_anchor="left",
        )

    _connect(
        slide.shapes,
        nodes["appcore"],
        nodes["autosave"],
        name="appcore_to_autosave",
        start_anchor="bottom",
        end_anchor="top",
    )
    _connect(
        slide.shapes,
        nodes["appcore"],
        nodes["logging"],
        name="appcore_to_logging",
        start_anchor="bottom",
        end_anchor="top",
    )
    _connect(
        slide.shapes,
        nodes["appcore"],
        nodes["reporter"],
        name="appcore_to_reporter",
        start_anchor="bottom",
        end_anchor="top",
    )
    _connect(
        slide.shapes,
        nodes["reporter"],
        nodes["files"],
        name="reporter_to_files",
        start_anchor="right",
        end_anchor="left",
    )

    return nodes



def _sequence_diagram(slide, *, content_top: float, slide_width: float) -> None:
    lifeline_height = PPTX_SLIDE_HEIGHT - content_top - PPTX_MARGIN - Cm(1)
    lifelines: dict[str, object] = {}

    participants = [
        ("gestion", "Gestión"),
        ("investigador", "Investigador"),
        ("ui", "UI Python Tkinter GUI\n(FormularioInvestigacionesFraude)"),
        ("valid", "Validaciones (validators.py)"),
        ("catalogos", "Servicio de Catálogos/Autopoblado"),
        ("importador", "Importador CSV"),
        ("autosave", "Autoguardado"),
        ("logs", "logger/logs.csv"),
        ("reporte", "Reportes (report_builder)"),
    ]

    interior_width = slide_width - 2 * PPTX_MARGIN
    spacing = interior_width / (len(participants) - 1)
    head_width = Cm(5)

    def lane_center(idx: float) -> float:
        return PPTX_MARGIN + spacing * idx

    for idx, (name, label) in enumerate(participants):
        left = PPTX_MARGIN + spacing * idx - head_width / 2
        head = _add_box(
            slide,
            name=f"lifeline_{name}",
            text=label,
            left=left,
            top=content_top,
            width=head_width,
            height=Cm(2.6),
            fill=PPTX_CONTEXT_FILL,
            bold=True,
            shape_type=MSO_SHAPE.RECTANGLE,
        )
        center_x = head.left + head.width / 2
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            center_x,
            head.top + head.height,
            center_x,
            content_top + lifeline_height,
        )
        line.name = f"lifeline_line_{name}"
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        line.line.width = Pt(1.5)
        lifelines[name] = head

    base_y = Cm(1)
    step_height = Cm(0.7)

    def message(step: int, source: str, target: str, text: str) -> None:
        start = lifelines[source]
        end = lifelines[target]
        start_x = start.left + start.width / 2
        end_x = end.left + end.width / 2
        y = content_top + base_y + step_height * (step - 1)
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, start_x, y, end_x, y
        )
        connector.name = f"msg_{step}_{source}_to_{target}"
        connector.line.width = Pt(2)
        text_box = slide.shapes.add_textbox(
            min(start_x, end_x),
            y - Cm(0.7),
            abs(end_x - start_x) + Cm(1.2),
            Cm(1.5),
        )
        text_box.name = f"label_{step}"
        para = text_box.text_frame.paragraphs[0]
        para.text = f"{step}. {text}"
        para.font.size = Pt(16)
        para.alignment = PP_ALIGN.LEFT

    # Mensajes
    message(1, "gestion", "investigador", "Asigna caso y solicita informe")
    message(2, "investigador", "ui", "Abre aplicación y selecciona pestañas")
    _add_note(
        slide,
        text="Entrada manual de datos",
        left=lane_center(2) - spacing * 1.5,
        top=content_top + base_y + step_height * 4,
        width=spacing * 3,
        height=Cm(1.4),
    )
    message(3, "investigador", "ui", "Escribe campos (Entry/Combobox/Text)")
    message(4, "ui", "valid", "FieldValidator valida formato/montos/fechas")
    message(5, "valid", "ui", "Errores inline o messagebox")
    message(6, "ui", "logs", "log_event(tipo=\"validacion/navegacion\")")
    message(7, "ui", "autosave", "request_autosave() (debounce 4s)")
    message(8, "autosave", "ui", "save_auto() + temp_version")

    _add_note(
        slide,
        text="Carga masiva",
        left=lane_center(2.5) - spacing * 1.1,
        top=content_top + base_y + step_height * 12,
        width=spacing * 2.2,
        height=Cm(1.4),
    )
    message(9, "investigador", "ui", "Acciones > Importar CSV")
    message(10, "ui", "importador", "_start_background_import(worker)")
    message(11, "importador", "ui", "payload de filas normalizadas")
    message(
        12,
        "ui",
        "ui",
        "_apply_*_import_payload pobla pestañas (Caso/Cliente/Equipo/Producto/Riesgo/Norma/Resumen)",
    )
    message(
        13,
        "ui",
        "valid",
        "Validación por fila y duplicados (cliente/colaborador/riesgo/producto/norma)",
    )
    message(14, "ui", "logs", "log_event(subtipo=\"import\")")
    message(15, "ui", "autosave", "_notify_dataset_changed()")

    _add_note(
        slide,
        text="Autocompletado/Autopoblado",
        left=lane_center(3.5) - spacing * 1.15,
        top=content_top + base_y + step_height * 20,
        width=spacing * 2.3,
        height=Cm(1.4),
    )
    message(
        16,
        "investigador",
        "ui",
        "Ingresa IDs (cliente/colaborador/producto/caso/riesgo/norma)",
    )
    message(17, "ui", "catalogos", "lookup en catálogos locales y snapshots")
    message(18, "catalogos", "ui", "Datos encontrados + metadatos")
    message(19, "ui", "valid", "should_autofill_field respeta campos editados")
    message(20, "ui", "ui", "Actualiza combobox/entries/tablas")
    message(
        21,
        "valid",
        "catalogos",
        "Validación de integridad referencial (IDs existen en catálogos/snapshots)",
    )
    message(22, "ui", "logs", "log_event(subtipo=\"autofill\")")

    _add_note(
        slide,
        text="Autosave y logging continuos",
        left=lane_center(1.5) - spacing * 1.5,
        top=content_top + base_y + step_height * 24,
        width=spacing * 3,
        height=Cm(1.4),
    )
    message(23, "ui", "autosave", "autosave_cycle() -> autosaves/<caso>/auto_N.json")
    message(24, "ui", "logs", "Encola evento y flush a logs.csv + external drive")

    _add_note(
        slide,
        text="Guardado final",
        left=lane_center(2.7) - spacing * 1.25,
        top=content_top + base_y + step_height * 27,
        width=spacing * 2.5,
        height=Cm(1.4),
    )
    message(25, "investigador", "ui", "Acciones > Guardar y enviar")
    message(26, "ui", "valid", "Validación global de todas las pestañas")
    message(27, "ui", "reporte", "build_report + build_docx + build_editable_deck")
    message(28, "reporte", "investigador", "Exporta JSON/CSV/Markdown/DOCX/PDF/PPTX")
    message(29, "ui", "autosave", "save_temp_version() + backup en external drive")
    message(30, "ui", "logs", "log_event(\"export\")")


def build_report(output: Path = DEFAULT_OUTPUT) -> Path:
    styles = _build_stylesheet()

    render_mermaid(ARCH_MMD, ARCH_PNG)
    _ensure_db_diagram_png()
    render_mermaid(SEQ_MMD, SEQ_PNG)

    page_margins = {
        "left": 0.7 * inch,
        "right": 0.7 * inch,
        "top": 0.8 * inch,
        "bottom": 0.8 * inch,
    }

    doc = BaseDocTemplate(
        str(output),
        pagesize=LETTER,
        leftMargin=page_margins["left"],
        rightMargin=page_margins["right"],
        topMargin=page_margins["top"],
        bottomMargin=page_margins["bottom"],
        title="Formulario de investigación – Arquitectura",
        author="Documento de referencia de prototipo",
    )
    frame = Frame(doc.leftMargin, doc.bottomMargin, doc.width, doc.height, id="normal")

    arch_page_size = landscape(A3)
    arch_frame = Frame(
        doc.leftMargin,
        doc.bottomMargin,
        arch_page_size[0] - doc.leftMargin - doc.rightMargin,
        arch_page_size[1] - doc.topMargin - doc.bottomMargin,
        id="arch_frame",
    )

    db_page_size = landscape(A3)
    db_frame = Frame(
        doc.leftMargin,
        doc.bottomMargin,
        db_page_size[0] - doc.leftMargin - doc.rightMargin,
        db_page_size[1] - doc.topMargin - doc.bottomMargin,
        id="db_frame",
    )

    seq_page_size = A3
    seq_frame = Frame(
        doc.leftMargin,
        doc.bottomMargin,
        seq_page_size[0] - doc.leftMargin - doc.rightMargin,
        seq_page_size[1] - doc.topMargin - doc.bottomMargin,
        id="seq_frame",
    )

    doc.addPageTemplates(
        [
            PageTemplate(id="main", frames=frame, onPage=_page_decor, pagesize=LETTER),
            PageTemplate(
                id="arch_diagram", frames=arch_frame, onPage=_page_decor, pagesize=arch_page_size
            ),
            PageTemplate(
                id="db_diagram", frames=db_frame, onPage=_page_decor, pagesize=db_page_size
            ),
            PageTemplate(
                id="seq_diagram", frames=seq_frame, onPage=_page_decor, pagesize=seq_page_size
            ),
        ]
    )
    doc.afterFlowable = lambda flowable: _after_flowable(doc, flowable)

    today = _dt.date.today().strftime("%Y-%m-%d")

    story = [
        Spacer(1, 1.5 * inch),
        Paragraph("Formulario de Investigación", styles["CoverTitle"]),
        Paragraph("Arquitectura y flujos de datos", styles["CoverSubtitle"]),
        Paragraph(f"Versión 1.0 — {today}", styles["Meta"]),
        Paragraph("Autor: Documentación referencia de prototipo", styles["Meta"]),
        Spacer(1, 1.2 * inch),
        Paragraph(
            "Documento generado automáticamente. Ejecuta build_architecture_report.py para actualizarlo.",
            styles["Meta"],
        ),
        PageBreak(),
    ]

    toc = TableOfContents()
    toc.levelStyles = [
        ParagraphStyle(name="TOCHeading1", fontSize=12, leftIndent=12, spaceAfter=4),
        ParagraphStyle(name="TOCHeading2", fontSize=10, leftIndent=24, spaceAfter=2),
    ]
    story.extend([_heading("Tabla de contenidos", styles, 1), Spacer(1, 0.2 * inch), toc, PageBreak()])

    # Section 1: Contexto
    story.append(_heading("1. Contexto del sistema y propósito", styles, 1))
    story.append(
        _paragraph(
            "Gestión asigna casos de fraude y los investigadores usan la aplicación Tkinter\n"
            "(main.py / app.py) para capturar datos de caso, clientes, colaboradores, productos,\n"
            "riesgos y normas. El objetivo es validar en línea, prevenir duplicados y producir\n"
            "un expediente exportable (CSV/JSON/Markdown/DOCX) desde report_builder.py.",
            styles,
        )
    )

    # Section 2: Flujos de captura
    story.append(_heading("2. Flujos de ingreso de datos", styles, 1))
    story.append(_heading("2.1 Ingreso manual", styles, 2))
    story.append(
        _paragraph(
            "Cada pestaña en ui/ frames utiliza ttk.Entry, Combobox y Text widgets con validación\n"
            "en blur (FieldValidator). Se verifican formatos de fechas, montos y claves para mantener\n"
            "consistencia con las reglas del Design document CM.",
            styles,
        )
    )
    story.append(_heading("2.2 Carga masiva", styles, 2))
    story.append(
        _paragraph(
            "Acciones > Importar procesa archivos .txt (pipes), .csv o .xlsx. El módulo utils.mass_load\n"
            "normaliza filas hacia el modelo interno y las pestañas de UI aplican _apply_*_import_payload\n"
            "para hidratar las tablas y disparar las mismas validaciones que el ingreso manual.",
            styles,
        )
    )
    story.append(_bullet_list(
        [
            "Formatos soportados: .txt con pipes, .csv, .xlsx",
            "Campos mapeados según utils/mass_load.py y cabeceras esperadas en cada entidad",
            "Filas inválidas se descartan con mensajes en la barra de estado y logs",
        ],
        styles,
    ))

    # Section 3: Entidades
    story.append(_heading("3. Entidades y campos", styles, 1))
    story.append(
        _paragraph(
            "Las pestañas Cliente, Equipo, Producto, Reclamo, Involucramiento, Riesgo y Norma usan\n"
            "combobox y entradas obligatorias. Los IDs siguen reglas: caso AAAA-NNNN, reclamo C########,\n"
            "team member letra+5 dígitos, analítica contable 10 dígitos (43/45/46/56) y validaciones\n"
            "de fechas y montos definidas en validators.py.",
            styles,
        )
    )

    # Section 4: Autocomplete
    story.append(_heading("4. Autocomplete y catálogos", styles, 1))
    story.append(
        _paragraph(
            "utils/autocomplete.py conecta Combobox con catálogos locales (client_details.csv,\n"
            "team_details.csv, risk_details.csv, norm_details.csv). En producción se espera consumir\n"
            "archivos descargados de Databricks para clientes/colaboradores, el registro GRC para riesgos\n"
            "y el índice maestro de normas.",
            styles,
        )
    )
    story.append(
        _bullet_list(
            [
                "Autocompletar preserva valores editados por el usuario (should_autofill_field).",
                "Búsquedas se disparan al perder foco o al confirmar el ID.",
                "Fuentes futuras: Databricks (clientes/colaboradores), GRC (riesgos), índice de normas, GDH mensual para activos y cesados.",
            ],
            styles,
        )
    )

    # Section 5: Autosave & logs
    story.append(_heading("5. Autosave y bitácora", styles, 1))
    story.append(
        _paragraph(
            "Autosave se ejecuta de forma diferida tras ediciones y en intervalos; genera autosave.json\n"
            "y autosaves/<timestamp>_<case_id>.json. Los eventos de validación, importación y navegación\n"
            "se almacenan en logs.csv para auditoría y para visualización con analytics/usage_visualizer.py.",
            styles,
        )
    )
    story.append(
        _bullet_list(
            [
                "Desencadenantes: pérdida de foco, temporizador y cambios masivos.",
                "Formato de log: timestamp, tipo, subtipo, widget_id, valor, case_id.",
                "Los respaldos también se copian a 'external drive/<case_id>/' cuando existe permiso.",
            ],
            styles,
        )
    )

    # Section 6: Generación de reportes
    story.append(_heading("6. Flujo de generación de reportes", styles, 1))
    story.append(
        _paragraph(
            "El botón 'Guardar y enviar' valida todas las pestañas (validators.py), construye el contexto con\n"
            "report_builder._build_report_context y exporta CSV, JSON, Markdown y DOCX. Si python-docx está\n"
            "instalado, se arma el informe Word usando la plantilla Plantilla_reporte.md como base.",
            styles,
        )
    )

    # Section 7: Tecnología
    story.append(_heading("7. Resumen de tecnología", styles, 1))
    story.append(_technology_table(styles))

    # Diagrams
    story.append(NextPageTemplate("arch_diagram"))
    story.append(PageBreak())
    story.append(_heading("Anexo A — Diagrama de arquitectura (Mermaid)", styles, 1))
    story.append(Spacer(1, 0.1 * inch))
    story.append(_flowable_image(ARCH_PNG, target_width=arch_frame.width))

    story.append(NextPageTemplate("db_diagram"))
    story.append(PageBreak())
    story.append(_heading("Anexo B — Modelo de datos para CSV", styles, 1))
    story.append(Spacer(1, 0.1 * inch))
    story.append(_flowable_image(DB_ARCH_PNG, target_width=db_frame.width))

    story.append(NextPageTemplate("seq_diagram"))
    story.append(PageBreak())
    story.append(_heading("Anexo C — Diagrama de secuencia", styles, 1))
    story.append(Spacer(1, 0.1 * inch))
    story.append(_flowable_image(SEQ_PNG, target_width=seq_frame.width))

    doc.build(story)
    return output


def build_editable_deck(output: Path = DEFAULT_PPTX) -> Path:
    """Generate a PowerPoint deck with editable architecture/sequence diagrams."""
    presentation = Presentation()
    presentation.slide_width = PPTX_SLIDE_WIDTH
    presentation.slide_height = PPTX_SLIDE_HEIGHT

    arch_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    arch_content_top = _add_header(
        arch_slide,
        title="Arquitectura – editable",
        subtitle=(
            "Diagrama construido con formas nativas para mover cajas, alinear conectores "
            "y ajustar estilo durante la revisión."
        ),
        slide_width=presentation.slide_width,
    )
    _architecture_nodes(
        arch_slide, content_top=arch_content_top, slide_width=presentation.slide_width
    )

    seq_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    seq_content_top = _add_header(
        seq_slide,
        title="Secuencia – editable",
        subtitle=(
            "Flujo de investigación, importación y exporte generado con conectores nativos para edición inmediata."
        ),
        slide_width=presentation.slide_width,
    )
    _sequence_diagram(
        seq_slide, content_top=seq_content_top, slide_width=presentation.slide_width
    )

    db_png = _ensure_db_diagram_png()
    db_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    db_content_top = _add_header(
        db_slide,
        title="Modelo de datos CSV",
        subtitle=(
            "Diagrama entidad-relación alineado a los CSV exportados (caso, clientes, "
            "colaboradores, productos, reclamos, riesgos, normas y llave técnica)."
        ),
        slide_width=presentation.slide_width,
    )
    max_width = presentation.slide_width - 2 * PPTX_MARGIN
    max_height = presentation.slide_height - db_content_top - PPTX_MARGIN
    img_width, img_height = _scale_image_to_box(db_png, max_width, max_height)
    db_left = PPTX_MARGIN + (max_width - img_width) / 2
    db_top = db_content_top + (max_height - img_height) / 2
    db_slide.shapes.add_picture(str(db_png), db_left, db_top, width=img_width, height=img_height)

    presentation.save(output)
    return output


def parse_args(args: list[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Genera el PDF de arquitectura y flujos.")
    parser.add_argument(
        "--output",
        type=Path,
        default=DEFAULT_OUTPUT,
        help=f"Ruta del PDF de salida (por defecto {DEFAULT_OUTPUT.name})",
    )
    parser.add_argument(
        "--pptx-output",
        type=Path,
        default=DEFAULT_PPTX,
        help=(
            "Ruta para generar la presentación editable con los diagramas. "
            "Se habilita por defecto; usa --no-pptx para omitirla"
        ),
    )
    parser.add_argument(
        "--no-pptx",
        action="store_true",
        help="No exportar la presentación editable.",
    )
    return parser.parse_args(args)


def main(argv: list[str] | None = None) -> int:
    argv = argv if argv is not None else sys.argv[1:]
    args = parse_args(argv)
    try:
        output = build_report(args.output)
    except Exception as exc:  # noqa: BLE001
        sys.stderr.write(f"Error generando el PDF: {exc}\n")
        return 1
    else:
        print(f"PDF generado en {output}")
    pptx_output: Path | None = None
    if not args.no_pptx and args.pptx_output:
        try:
            pptx_output = build_editable_deck(args.pptx_output)
        except Exception as exc:  # noqa: BLE001
            sys.stderr.write(f"Error generando la presentación editable: {exc}\n")
            return 1

    if pptx_output:
        print(f"Presentación editable generada en {pptx_output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
